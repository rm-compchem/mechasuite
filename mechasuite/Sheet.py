import sys
import os
from PyQt5.QtWidgets import QApplication, QMainWindow, QAction, QWidget, \
    QMenu, QTableWidget, QTableWidgetItem, QListView, QHBoxLayout, QVBoxLayout, QSplitter, \
    QDialog, QPushButton, QFileDialog, QHeaderView, QInputDialog, QMessageBox, \
    QFrame, QOpenGLWidget, QGridLayout, QComboBox, QTextEdit, QAbstractItemView, QPushButton, \
    QListWidget, QDialogButtonBox, QListWidgetItem
from PyQt5.QtGui import QPixmap, QCursor, QFont, QColor
from PyQt5.QtCore import pyqtSlot, pyqtSignal
from .MyWidgets import *
from .Plot import plot, MainCanvas
import json
from OpenGL.GL import *
from OpenGL.GLU import *
from .Data import *
from math import cos, sin, acos
from .Read import unfold_temps
import re
# import rotate
import time

class Log:
    def __init__(self):
        return

    @staticmethod
    def show_msg(msg: str):
        msgb = QMessageBox()
        msgb.setText(msg)
        msgb.exec()

class InfoLabel(QTextEdit):
    textchanged = pyqtSignal(str)

    def __init__(self, parent):
        QTextEdit.__init__(self, parent)
        self.setPlainText("")
        # self.setStyleSheet("background-color:gray;")
        # self.info_label.setWordWrap(True)
        self.setMaximumHeight(120)
        self.setFrameStyle(QFrame.Box)

    def keyPressEvent(self, event):
        self.textchanged.emit(self.toPlainText())
        QTextEdit.keyPressEvent(self, event)


class QGLBegin(QOpenGLWidget):
    def __init__(self, parent):
        QOpenGLWidget.__init__(self, parent)

        #self.setMaximumSize(800, 800)
        self.setMinimumSize(200, 200)

        self.setMouseTracking(True)

        self.rRange, self.rRangeVP = 10.0, 10.0
        self.lRange, self.lRangeVP = -10.0, -10.0
        self.uRange, self.uRangeVP = 10.0, 10.0
        self.dRange, self.dRangeVP = -10.0, -10.0

        self.coors = []
        self.labels = []
        self.cell = [[0, 0, 0], [0, 0, 0], [0, 0, 0]]
        self.cell_lines = [[0, 0, 0] for i in range(24)]
        self.showcell = False
        self.bonds = []
        self.mouseX, self.mouseY, self.mouseYbef, self.mouseXbef = 200, 200, 200, 200
        self.mouseXbef, self.mouseYbef = 0, 0
        self.dx, self.dy = 0, 0
        self.xrot, self.yrot, self.rot_ang = 0.0, 0.0, 0.02
        self.rot_mat = [[1, 0, 0], [0, 1, 0], [0, 0, 1]]
        self.struct = None
        self.bool_peli = False
        self.struct_is_set = False
        self.zlimit = 0
        self.initializeGL()

    def set_zrange(self):
        # return
        if len(self.coors) < 50:
            self.zlimit = 0
            return
        tempcapas = [self.coors[0][2]]
        for coor in self.coors:
            found = False
            for tc in tempcapas:
                # print(abs(coor[2] - tc))
                if abs(coor[2] - tc) < 0.6:
                    found = True
                    break
            if not found:
                tempcapas.append(coor[2])
        tempcapas.sort()
        if len(tempcapas)>5:
            self.zlimit = tempcapas[int(len(tempcapas)/10)]
        else:
            self.zlimit = 0

    def update_data(self, struct):
        # self.showcell = self.struct.showcell
        #if not struct.nat or not struct.labels:
        #    self.struct_is_set = False
        #    self.update()
        #    self.hide()
        #    return
        self.show()
        self.struct_is_set = True
        self.struct = struct
        self.lRange = struct.lRange
        self.rRange = struct.rRange
        self.dRange = struct.dRange
        self.uRange = struct.uRange
        self.resizeGL(self.width(),self.height())
        if struct.coors:
            try:
                self.coors = np.dot(struct.coors, struct.rot_matrix)
            except:
                self.coors = []
                self.struct_is_set = False
        else:
            self.coors = []
            self.struct_is_set = False
        self.labels = struct.labels
        # ni siquiera necesito tener self.cell lo dejare temporalemente
        self.cell = struct.cell[:]
        if self.cell[0][0] or self.cell[0][1] or self.cell[0][2]:
            self.set_cell_lines()
            self.cell_lines = np.dot(self.cell_lines, struct.rot_matrix)
        if len(self.coors) > 0:
            self.shift_to_mass_center()
        self.genbonds()
        self.update()

    def initializeGL(self):
        glClear(GL_COLOR_BUFFER_BIT | GL_DEPTH_BUFFER_BIT)
        glEnable(GL_BLEND)
        glBlendFunc(GL_SRC_ALPHA, GL_ONE_MINUS_SRC_ALPHA)
        self.enableGlobProp()

    def paintGL(self):
        glClearColor(0.0, 0.0, 0.0, 1.0)
        glClear(GL_COLOR_BUFFER_BIT | GL_DEPTH_BUFFER_BIT);
        glEnable(GL_BLEND);
        glBlendFunc(GL_SRC_ALPHA, GL_ONE_MINUS_SRC_ALPHA);

        self.render_scene()

    def resizeGL(self, w: int, h: int):
        rRange = self.rRange
        lRange = self.lRange
        uRange = self.uRange
        dRange = self.dRange
        nRange = 15.0
        self.makeCurrent();
        # glGetError()
        if h == 0: h = 1;

        glViewport(0, 0, w, h);
        glMatrixMode(GL_PROJECTION);
        glLoadIdentity();

        if w <= h:
            self.lRangeVP = lRange;
            self.rRangeVP = rRange;
            self.dRangeVP = dRange * h / w;
            self.uRangeVP = uRange * h / w;
            # // gluPerspective(90.0, -nRange * self.height / self.width, -nRange, nRange);
            glOrtho(lRange, rRange, dRange * h / w, uRange * h / w, -4 * rRange, 4 * rRange);
        else:
            self.lRangeVP = lRange * w / h;
            self.rRangeVP = rRange * w / h;
            self.dRangeVP = dRange;
            self.uRangeVP = uRange;
            # // gluPerspective(90.0, -nRange * self.width / self.height, -nRange, nRange);
            glOrtho(lRange * w / h, rRange * w / h, dRange, uRange, -4 * rRange, 4 * rRange);

        glMatrixMode(GL_MODELVIEW);
        glLoadIdentity();

    def enableGlobProp(self):
        lightPos = [-50.0, 50.0, 100.0, 1.0]
        ambientLight = [0.35, 0.35, 0.35, 1]
        diffuseLight = [0.75, 0.75, 0.75, 0.7]
        specular = [1.0, 1.0, 1.0, 1.0]
        specref = [0.0, 0.0, 0.0, 1.0]
        shininess = 128
        glLightfv(GL_LIGHT0, GL_POSITION, lightPos);
        glEnable(GL_DEPTH_TEST);
        glEnable(GL_LIGHTING);
        glLightfv(GL_LIGHT0, GL_AMBIENT, ambientLight);
        glLightfv(GL_LIGHT0, GL_DIFFUSE, diffuseLight);
        glLightfv(GL_LIGHT0, GL_SPECULAR, specular);
        glEnable(GL_LIGHT0);
        glEnable(GL_COLOR_MATERIAL);
        glColorMaterial(GL_FRONT, GL_AMBIENT_AND_DIFFUSE);
        glMaterialfv(GL_FRONT, GL_SPECULAR, specref);
        glMateriali(GL_FRONT, GL_SHININESS, shininess);
        glColor3ub(230, 100, 100);

    def genbonds(self):
        self.bonds = []
        for i, coor in enumerate(self.coors):
            self.bonds.append([[], []])
            for j in range(i):
                # calculate the distance between the two atoms self.coors[i] y self.coors[j]
                dist = np.array(coor) - np.array(self.coors[j])
                dist = dist ** 2
                dist = np.sqrt(dist.sum())
                if dist <= radios[self.labels[i]] + radios[self.labels[j]]:
                    self.bonds[-1][0].append(j)
                    if cylinders[self.labels[i]] <= cylinders[self.labels[j]]:
                        self.bonds[-1][1].append(cylinders[self.labels[i]])
                    else:
                        self.bonds[-1][1].append(cylinders[self.labels[j]])

    def set_cell_lines(self):
        self.cell_lines[0] = [0.0, 0.0, 0.0];
        self.cell_lines[1] = [self.cell[0][0], self.cell[0][1], self.cell[0][2]];
        # // de cero a b
        self.cell_lines[2] = [0.0, 0.0, 0.0];
        self.cell_lines[3] = [self.cell[1][0], self.cell[1][1], self.cell[1][2]];
        # // de cero a c
        self.cell_lines[4] = [0.0, 0.0, 0.0];
        self.cell_lines[5] = [self.cell[2][0], self.cell[2][1], self.cell[2][2]];
        # // de a a a+b
        self.cell_lines[6] = [self.cell[0][0], self.cell[0][1], self.cell[0][2]];
        self.cell_lines[7] = [self.cell[0][0] + self.cell[1][0], self.cell[0][1] + self.cell[1][1],
                              self.cell[0][2] + self.cell[1][2]];
        # // de  b a a+b
        self.cell_lines[8] = [self.cell[1][0], self.cell[1][1], self.cell[1][2]];
        self.cell_lines[9] = [self.cell[0][0] + self.cell[1][0], self.cell[0][1] + self.cell[1][1],
                              self.cell[0][2] + self.cell[1][2]];
        # // de a a a+c
        self.cell_lines[10] = [self.cell[0][0], self.cell[0][1], self.cell[0][2]];
        self.cell_lines[11] = [self.cell[0][0] + self.cell[2][0], self.cell[0][1] + self.cell[2][1],
                               self.cell[0][2] + self.cell[2][2]];
        # // de c a a+c
        self.cell_lines[12] = [self.cell[2][0], self.cell[2][1], self.cell[2][2]];
        self.cell_lines[13] = [self.cell[0][0] + self.cell[2][0], self.cell[0][1] + self.cell[2][1],
                               self.cell[0][2] + self.cell[2][2]];
        # // de  b a b+c
        self.cell_lines[14] = [self.cell[1][0], self.cell[1][1], self.cell[1][2]];
        self.cell_lines[15] = [self.cell[1][0] + self.cell[2][0], self.cell[1][1] + self.cell[2][1],
                               self.cell[1][2] + self.cell[2][2]];
        # // de  c a b+c
        self.cell_lines[16] = [self.cell[2][0], self.cell[2][1], self.cell[2][2]];
        self.cell_lines[17] = [self.cell[1][0] + self.cell[2][0], self.cell[1][1] + self.cell[2][1],
                               self.cell[1][2] + self.cell[2][2]];
        # //de b+c a a+b+c
        self.cell_lines[18] = [self.cell[1][0] + self.cell[2][0], self.cell[1][1] + self.cell[2][1],
                               self.cell[1][2] + self.cell[2][2]];
        self.cell_lines[19] = [self.cell[0][0] + self.cell[1][0] + self.cell[2][0],
                               self.cell[0][1] + self.cell[1][1] + self.cell[2][1],
                               self.cell[0][2] + self.cell[1][2] + self.cell[2][2]];
        # //de a+c a a+b+c
        self.cell_lines[20] = [self.cell[0][0] + self.cell[2][0], self.cell[0][1] + self.cell[2][1],
                               self.cell[0][2] + self.cell[2][2]];
        self.cell_lines[21] = [self.cell[0][0] + self.cell[1][0] + self.cell[2][0],
                               self.cell[0][1] + self.cell[1][1] + self.cell[2][1],
                               self.cell[0][2] + self.cell[1][2] + self.cell[2][2]];
        # //de a+b a a+b+c
        self.cell_lines[22] = [self.cell[0][0] + self.cell[1][0], self.cell[0][1] + self.cell[1][1],
                               self.cell[0][2] + self.cell[1][2]];
        self.cell_lines[23] = [self.cell[0][0] + self.cell[1][0] + self.cell[2][0],
                               self.cell[0][1] + self.cell[1][1] + self.cell[2][1],
                               self.cell[0][2] + self.cell[1][2] + self.cell[2][2]];

    def shift_to_mass_center(self):
        Xcm, Ycm, Zcm = 0.0, 0.0, 0.0
        MM = 0
        for i, coor in enumerate(self.coors):
            x = coor[0];
            y = coor[1];
            z = coor[2];
            Xcm += masas[self.labels[i]] * x
            Ycm += masas[self.labels[i]] * y
            Zcm += masas[self.labels[i]] * z
            MM += masas[self.labels[i]]

        Xcm /= MM;
        Ycm /= MM;
        Zcm /= MM;

        for i, _ in enumerate(self.coors):
            self.coors[i][0] -= Xcm;
            self.coors[i][1] -= Ycm;
            self.coors[i][2] -= Zcm;
        for i, _ in enumerate(self.cell_lines):
            self.cell_lines[i][0] -= Xcm
            self.cell_lines[i][1] -= Ycm
            self.cell_lines[i][2] -= Zcm

    def render_bonds(self, coor1, coor2, radio, subdiv, color):
        x1 = coor1[0];
        y1 = coor1[1];
        z1 = coor1[2];
        x2 = coor2[0];
        y2 = coor2[1];
        z2 = coor2[2];

        vx = x2 - x1;
        vy = y2 - y1;
        vz = z2 - z1;

        if (vz == 0): vz = .0001;

        v = np.sqrt(vx * vx + vy * vy + vz * vz);
        ax = 57.2957795 * acos(vz / v);
        if (vz < 0.0): ax = -ax;
        rx = -vy * vz;
        ry = vx * vz;
        glPushMatrix();

        glTranslatef(x1, y1, z1);
        glRotatef(ax, rx, ry, 0.0);
        glColor3f(color[0], color[1], color[2]);
        quadric = gluNewQuadric()
        gluQuadricOrientation(quadric, GLU_OUTSIDE);
        gluCylinder(quadric, radio, radio, v, subdiv, 1);
        glPopMatrix();

    def render_scene(self):
        self.set_zrange()
        for i, coor in enumerate(self.coors):
            if coor[0] < self.lRangeVP or coor[0] > self.rRangeVP or coor[1] < self.dRangeVP or coor[1] > self.uRangeVP:
                continue
            if not self.zlimit or coor[2] >= self.zlimit:
                for j, at2 in enumerate(self.bonds[i][0]):
                    # // at1 = i;
                    # at2 = bonds[i].neighbors[j];
                    semix = (coor[0] + self.coors[at2][0]) / 2.0
                    semiy = (coor[1] + self.coors[at2][1]) / 2.0
                    semiz = (coor[2] + self.coors[at2][2]) / 2.0
                    semis = (semix, semiy, semiz)

                    quadric = gluNewQuadric()
                    gluQuadricNormals(quadric, GLU_SMOOTH)
                    self.render_bonds(coor, semis, self.bonds[i][1][j], 24, colors[self.labels[i]])
                    gluDeleteQuadric(quadric)

                    quadric = gluNewQuadric()
                    gluQuadricNormals(quadric, GLU_SMOOTH)
                    self.render_bonds(semis, self.coors[at2], self.bonds[i][1][j], 24, colors[self.labels[at2]])
                    gluDeleteQuadric(quadric)

            glPushMatrix()
            glTranslatef(coor[0], coor[1], coor[2])
            glColor4f(colors[self.labels[i]][0], colors[self.labels[i]][1], colors[self.labels[i]][2],
                      colors[self.labels[i]][3])
            # glColor4f(coor_colors[i][0], coor_colors[i][1], coor_colors[i][2], coor_colors[i][3]);
            quadric = gluNewQuadric();
            gluQuadricNormals(quadric, GLU_SMOOTH);
            # gluSphere(quadric, coor_spheres[i], quality, 2);
            gluSphere(quadric, spheres[self.labels[i]], 24, 2);
            gluDeleteQuadric(quadric);
            glPopMatrix();

            # // cout << "tres\n";
            # for (int i=0;i < selected_atoms.size();i++){
            # s_a_index = selected_atoms[i];
            # glPushMatrix();
            # glEnable(GL_BLEND);
            # glBlendFunc(GL_SRC_ALPHA, GL_ONE_MINUS_SRC_ALPHA);
            # glTranslatef((GLfloat)coors[s_a_index][0], (GLfloat)coors[s_a_index][1], (GLfloat)coors[s_a_index][2]);
            # glColor4f(0.0, 0.0, 0.9, 0.3);
            # // glutSolidSphere(coor_spheres[s_a_index]+0.2, 15, 10);
            # quadric=gluNewQuadric();
            # gluQuadricNormals(quadric, GLU_SMOOTH);
            # gluSphere(quadric, coor_spheres[s_a_index]+0.2, 20, 2);
            # gluDeleteQuadric(quadric);
            # glPopMatrix();
            #
            # }
            #
            # if (paint_sel_rectangle){
            # glColor4f(0.0, 0.0, 0.5, 0.3);
            # // glEnable(GL_BLEND);
            # // glBlendFunc(GL_SRC_ALPHA, GL_ONE_MINUS_SRC_ALPHA);
            # glBegin(GL_QUADS);
            # glVertex3f(selrectx, selrecty, 3 * rRange);
            # glVertex3f(selrectx, mouseY, 3 * rRange);
            # glVertex3f(mouseX, mouseY, 3 * rRange);
            # glVertex3f(mouseX, selrecty, 3 * rRange);
            # glEnd();
            # }

        if self.showcell:
            if self.cell[0][0] or self.cell[0][1] or self.cell[0][2]:
                glColor4f(1.0, 1.0, 1.0, 0.8)
                # glEnable(GL_LINE_WIDTH)
                # glLineWidth(1)
                glBegin(GL_LINES)
                for cell_line in self.cell_lines:
                    glVertex3f(cell_line[0], cell_line[1], cell_line[2])

                glEnd()

    def mouseMoveEvent(self, e):
        self.mouseXbef = self.mouseX
        self.mouseYbef = self.mouseY
        self.mouseX = e.x()
        self.mouseY = e.y()
        self.dx = (self.mouseX - self.mouseXbef)
        self.dy = (self.mouseY - self.mouseYbef)

        if e.buttons() == Qt.LeftButton:
            self.rotate_coors()
            self.update()
        if e.buttons() == Qt.RightButton:
            self.setCursor(Qt.ClosedHandCursor)
            if self.dy != 0:
                self.dy = self.dy / self.height()
                self.dy *= (self.uRangeVP-self.dRangeVP)
                self.dRange += self.dy
                self.uRange += self.dy

            if self.dx != 0:
                self.dx /= self.width()
                self.dx *= (self.rRangeVP - self.lRangeVP)
                self.rRange -= self.dx
                self.lRange -= self.dx
            self.resizeGL(self.width(), self.height())
            self.update()
            self.struct.lRange = self.lRange
            self.struct.rRange = self.rRange
            self.struct.dRange = self.dRange
            self.struct.uRange = self.uRange

    def mouseReleaseEvent(self, e):
        self.setCursor(Qt.ArrowCursor)

    def wheelEvent(self, e):
        # print(e.angleDelta().y())
        if e.angleDelta().y() > 0:
            if self.rRange - self.lRange <= 2:
                return 
            if self.uRange - self.dRange <= 2:
                return
            self.rRange -= 1
            self.lRange += 1
            self.uRange -= 1
            self.dRange += 1
        else:
            self.rRange += 1
            self.lRange -= 1
            self.uRange += 1
            self.dRange -= 1

        self.resizeGL(self.width(), self.height())
        self.update()
        self.struct.lRange = self.lRange
        self.struct.rRange = self.rRange
        self.struct.dRange = self.dRange
        self.struct.uRange = self.uRange

    def mouseDoubleClickEvent(self, evt):
        return
        if self.bool_peli:
            self.bool_peli = False
        else:
            self.bool_peli = True

        if self.bool_peli:
            for i in [-1, -1, -1, -1, -5, 1, 1, 1, 1, 1]:
                time.sleep(0.5)
                self.coors[0] += i
                # print(self.coors[0])
                self.render_scene()
                self.paintGL()
                self.update()

                print("sleeping")

    def rotate_coors(self):
        if not self.struct_is_set:
            return
        vector = np.array([self.dy, self.dx, 0.0])
        norm = np.linalg.norm(vector)
        if norm == 0:
            return
        uvector = vector / norm
        x, y, z = uvector
        ang = -self.rot_ang
        # multiply angle to make it real(faster)
        ang *= norm

        matrix = np.array([
            [cos(ang) + pow(x, 2) * (1 - cos(ang)),
             x * y * (1 - cos(ang)) - z * sin(ang),
             x * z * (1 - cos(ang)) + y * sin(ang)],
            [x * y * (1 - cos(ang)) + z * sin(ang),
             cos(ang) + pow(y, 2) * (1 - cos(ang)),
             y * z * (1 - cos(ang)) - x * sin(ang)],
            [z * x * (1 - cos(ang)) - y * sin(ang),
             z * y * (1 - cos(ang)) + x * sin(ang),
             cos(ang) + pow(z, 2) * (1 - cos(ang))]
        ]
        )

        # for i, coors in enumerate(self.coors):
        #     self.coors[i] = np.matmul(coors, matrix)
        self.coors = np.dot(self.coors, matrix)
        self.cell_lines = np.dot(self.cell_lines, matrix)

        if self.struct is not None:
            self.struct.rot_matrix = np.dot(self.struct.rot_matrix, matrix)

    def rotate_coors2(self):
        if self.struct is not None:
            self.coors, self.cell_lines, self.rot_mat = rotate.rotate_coors(self.dy, self.dx, self.rot_ang, self.coors, self.cell_lines, self.rot_mat)
        else:
            self.coors, self.cell_lines, _ = rotate.rotate_coors(self.dy, self.dx, self.rot_ang, self.coors,
                                                                            self.cell_lines)


class RefDialog(QDialog):
    def __init__(self, mec):
        QDialog.__init__(self)
        self.items = []
        self.names = mec
        self.ok = False
        self.setWindowTitle("Zeros")

        # aqui quiero poner un input control para poner el nombre
        # y no tener que pedirlo de antemano

        self.list = QListView(self)
        self.model = QStandardItemModel(self.list)
        self.list.doubleClicked.connect(self.set_items)

        self.label = QLabel("")

        for name in self.names:
            item = QStandardItem(name)
            item.setEditable(False)
            # item.setCheckable(True)
            self.model.appendRow(item)
        self.list.setModel(self.model)
        self.list.show()

        self.okbutton = QPushButton("Ok")
        self.okbutton.clicked.connect(self.set_ok)
        self.closebutton = QPushButton("Cancel")
        self.closebutton.clicked.connect(self.unset_ok)
        self.resetbutton = QPushButton("Reset")
        self.resetbutton.clicked.connect(self.reset)

        self.layout = QVBoxLayout()
        self.sublayout = QHBoxLayout()
        self.layout.addWidget(self.list)
        self.layout.addWidget(self.label)
        self.layout.addLayout(self.sublayout)
        self.sublayout.addWidget(self.okbutton)
        self.sublayout.addWidget(self.closebutton)
        self.sublayout.addWidget(self.resetbutton)
        self.setLayout(self.layout)

        self.exec()

    def reset(self, e):
        self.items = []
        self.label.setText("")

    def set_ok(self, e):
        self.ok = True
        self.close()

    def unset_ok(self, e):
        self.ok = False
        self.items = []
        self.close()

    def set_items(self, item):
        row = item.row()
        self.items.append(self.names[row])
        self.label.setText("+".join(self.items))


class ThermoDialog(QDialog):
    def __init__(self, parent):
        QDialog.__init__(self, parent)

        self.temps = []
        self.modes = 0
        self.P = 1
        self.V = 0
        self.ok = False
        self.no_neg_freq = False

        self.labelT = QLabel("T (K):")
        self.labelModes = QLabel("Normal Modes:")
        self.labelP = QLabel("Pressure (atm): ")
        self.labelV = QLabel("Volume (L): ")
        self.checkbox = QCheckBox("Ignore negative frequencies")
        self.lineEditT = QLineEdit()
        self.lineEditModes = QLineEdit()
        self.lineEditP = QLineEdit()
        self.lineEditP.setText("1")
        self.lineEditV = QLineEdit()
        self.lineEditV.setText("0")
        self.okb = QPushButton("Ok")
        self.cancelb = QPushButton("Cancel")

        self.okb.clicked.connect(self.on_accept)
        self.cancelb.clicked.connect(self.reject)

        self.glayout = QGridLayout()
        self.glayout.addWidget(self.labelT, 0, 0)
        self.glayout.addWidget(self.lineEditT, 0, 1)
        self.glayout.addWidget(self.labelP, 1, 0)
        self.glayout.addWidget(self.lineEditP, 1, 1)
        self.glayout.addWidget(self.labelV, 2, 0)
        self.glayout.addWidget(self.lineEditV, 2, 1)
        self.glayout.addWidget(self.labelModes, 3, 0)
        self.glayout.addWidget(self.lineEditModes, 3, 1)
        self.glayout.addWidget(self.checkbox, 4, 1)
        self.glayout.addWidget(self.okb, 5, 0)
        self.glayout.addWidget(self.cancelb, 5, 1)
        self.setLayout(self.glayout)

    def get_values(self):
        self.exec()
        return self.ok, self.temps, self.modes, self.P, self.V, self.no_neg_freq  # self.pg

    def on_accept(self):
        self.temps = unfold_temps(self.lineEditT.text())
        if not self.temps:
            self.lineEditT.setText("Enter valid T")
            return
        self.modes = self.lineEditModes.text()
        if self.modes == "":
            self.modes = 0
        else:
            try:
                self.modes = int(self.modes)
            except ValueError:
                self.lineEditModes.setText("Enter valid value")
                return

        try:
            self.P = float(self.lineEditP.text())
        except ValueError:
            self.P = 1.0
            self.lineEditP.setText("Enter valid value")
            return

        try:
            self.V = float(self.lineEditV.text())
        except ValueError:
            self.V = 0.0
            self.lineEditV.setText("Enter valid value or zero")
            return

        self.ok = True
        self.no_neg_freq = self.checkbox.checkState()
        self.accept()


class PGDialog(QDialog):
    def __init__(self, parent, name):
        QDialog.__init__(self, parent)

        self.pg = ""
        self.ok = False
        self.labelPG = QLabel("Point group\n for "+name+":")

        self.comboPG = QComboBox()
        self.comboPG.addItems(["solid"] + list(Thermo.sigma.keys()))
        self.okb = QPushButton("Ok")
        self.cancelb = QPushButton("Cancel")

        self.okb.clicked.connect(self.on_accept)
        self.cancelb.clicked.connect(self.reject)

        self.glayout = QGridLayout()
        self.glayout.addWidget(self.labelPG, 0, 0)
        self.glayout.addWidget(self.comboPG, 0, 1)
        self.glayout.addWidget(self.okb, 1, 0)
        self.glayout.addWidget(self.cancelb, 1, 1)
        self.setLayout(self.glayout)

    def get_values(self):
        self.exec()
        return self.ok, self.pg

    def on_accept(self):
        self.pg = self.comboPG.currentText()
        self.ok = True
        self.accept()

class MECPropsDialog(QDialog):
    def __init__(self, parent=None, redmass=None, SOC=None, grad=None, diffgrad=None):
        super().__init__(parent)
        self.setWindowTitle("Set MECP Properties")
        self.redmass_edit = QLineEdit(self)
        self.SOC_edit = QLineEdit(self)
        self.grad_edit = QLineEdit(self)
        self.diffgrad_edit = QLineEdit(self)

        if redmass is not None:
            self.redmass_edit.setText(str(redmass))
        if SOC is not None:
            self.SOC_edit.setText(str(SOC))
        if grad is not None:
            self.grad_edit.setText(str(grad))
        if diffgrad is not None:
            self.diffgrad_edit.setText(str(diffgrad))

        self.okb = QPushButton("Ok")
        self.cancelb = QPushButton("Cancel")
        self.okb.clicked.connect(self.accept)
        self.cancelb.clicked.connect(self.reject)

        layout = QGridLayout()
        layout.addWidget(QLabel("RedMass (a.m.u.):"), 0, 0)
        layout.addWidget(self.redmass_edit, 0, 1)
        layout.addWidget(QLabel("SOC (cm-1):"), 1, 0)
        layout.addWidget(self.SOC_edit, 1, 1)
        layout.addWidget(QLabel("Gradient (Ha/B):"), 2, 0)
        layout.addWidget(self.grad_edit, 2, 1)
        layout.addWidget(QLabel("DiffGrad (Ha/B):"), 3, 0)
        layout.addWidget(self.diffgrad_edit, 3, 1)
        layout.addWidget(self.okb, 4, 0)
        layout.addWidget(self.cancelb, 4, 1)
        self.setLayout(layout)

    def get_values(self):
        if self.exec() != QDialog.Accepted:
            return False, None, None, None, None
        try:
            redmass = float(self.redmass_edit.text())
            SOC = float(self.SOC_edit.text())
            grad = float(self.grad_edit.text())
            diffgrad = float(self.diffgrad_edit.text())
        except ValueError:
            msg = QMessageBox(self)
            msg.setText("All values must be numeric.")
            msg.exec()
            return False, None, None, None, None
        return True, redmass, SOC, grad, diffgrad


class MainSheet(QTableWidget):
    def __init__(self, parent, data=None, rows=50, cols=20):
        QTableWidget.__init__(self)
        self.setRowCount(rows)
        self.setColumnCount(cols)
        self.move(0, 0)
        self.parent = parent
        self.data = data
        self.activecelltext = ""
        self.editing = False
        self.setGridStyle(Qt.DotLine)

        # bindings-----------------------------------
        # table selection change
        # self.doubleClicked.connect(self.on_dclick)
        self.itemClicked.connect(self.on_itm_click)
        self.itemDoubleClicked.connect(self.on_itm_dclick)
        # self.cellDoubleClicked.connect(self.on_cell_dclick)
        # self.cellClicked.connect(self.on_cell_click)
        self.itemChanged.connect(self.itm_changed)
        self.currentItemChanged.connect(self.currentitemchanged)
        # bindings-----------------------------------

        headers = self.horizontalHeader()
        headers.setSectionsClickable(True)
        # headers.sectionEntered.connect(self.columnclicked)
        headers.setContextMenuPolicy(Qt.CustomContextMenu)
        headers.customContextMenuRequested.connect(self.header_popup)

    # def columnclicked(self, a):
    #     print("sii", a)

    def currentitemchanged(self, itm, itm2):
        self.editing = False
        if itm is None:
            return
        self.activecelltext = itm.text()
        self.show_itm_info(itm)

    @pyqtSlot()
    def keyPressEvent(self, itm):
        key = itm.key()
        self.editing = True
        if key == 16777223:
            self.editing = False
            self.del_itms()
        if key == 16777216:
            self.editing = False
        # if key == Qt.Key_F2:
        #     self.editing = True
        #     print("editing")
        QTableWidget.keyPressEvent(self, itm)

    def itm_changed(self, itm):
        if not self.editing:
            return
        itmtext, row, col = itm.text(), itm.row(), itm.column()
        if itmtext == self.activecelltext:
            return

        msg = QMessageBox()
        msg.setIcon(QMessageBox.Information)

        cell_item_obj = self.item(row - 1, col)
        # should not try to change non existint items
        if cell_item_obj is None:
            msg.setText("Item does not exist")
            msg.setWindowTitle("Error")
            msg.exec()
            itm.setText("")
            return

        # print("item changed")
        rowname = self.verticalHeaderItem(row).text()
        colname = self.horizontalHeaderItem(col).text()
        if rowname == "Name":
            colobj = self.data.get_mech(colname)
            if colobj.itm_exists(itmtext):
                msg.setText(itmtext + " already exists")
                msg.setWindowTitle("Error")
                msg.exec()
                itm.setText(self.activecelltext)
            else:
                colobj.itm_change_name(self.activecelltext, itmtext)
        elif rowname == "Energy":
            colobj = self.data.get_mech(colname)

            itmname = cell_item_obj.text()

            if colobj.get_itm(itmname) is None or not colobj.get_itm(itmname).set_e(itmtext):
                msg.setText(itmtext + " bad energy value")
                msg.setWindowTitle("Error")
                msg.exec()
                itm.setText(self.activecelltext)
        else:
            msg.setText("Cell cannot be modified!")
            msg.setWindowTitle("Error")
            msg.exec()
            itm.setText("")

    def header_popup(self, event):
        # self.horizontalHeader().sectionClicked.emit(1)
        menu = QMenu(self)
        # setdefbasisref = menu.addAction("Set default basis ref")
        newitm = menu.addAction("New Item")
        order_itms = menu.addAction("Order Items")
        mergeItems = menu.addAction("Merge Items")
        importItm = menu.addAction("Import Intermediate")
        addref = menu.addAction("Add Reference")
        chgunit = menu.addAction("Convert Unit")
        chgname = menu.addAction("Change Mechanism Name")
        delgraph = menu.addAction("Delete Plot")
        del_refs = menu.addAction("Delete References")
        action = menu.exec_(self.mapToGlobal(event))
        if action == importItm:
            self.import_itm()
        # elif action == setdefbasisref:
        #     self.set_def_basis_ref()
        elif action == addref:
            self.add_ref_to_mech()
        elif action == mergeItems:
            self.on_header_merge_items()
        elif action == chgunit:
            self.chg_unit()
        elif action == newitm:
            self.new_item()
        elif action == delgraph:
            self.del_graph()
        elif action == del_refs:
            self.del_refs_from_mec()
        elif action == chgname:
            self.change_mec_name()
        elif action == order_itms:
            self.order_itms()

    def order_itms(self):
        col = self.horizontalHeaderItem(self.currentColumn()).text()
        itm_name = self.data.get_mech(col).get_itms_names()

        ow = OrderWindow(self, itm_name, col)

    def change_mec_name(self):
        newname, ok = QInputDialog.getText(self, "Enter new name", "New name")
        if not ok:
            return
        col = self.currentColumn()
        colname = self.horizontalHeaderItem(col).text()
        mec = self.data.change_mec_name(colname, newname)
        self.update_data()

    def del_refs_from_mec(self):
        col = self.currentColumn()
        colname = self.horizontalHeaderItem(col).text()
        mec = self.data.get_mech(colname)
        refsnames = mec.get_refs_names()
        dialog = MultipleChoiceDialog(refsnames+["Default Reference"])
        dialog.exec()
        if not dialog.ok:
            return
        for refname in dialog.items:
            if refname == "Default Reference":
                refname = "Def Ref"
            mec.del_ref(refname)

    def del_graph(self):
        pnames = self.data.get_plots_names()
        dialog = MultipleChoiceDialog(pnames)
        dialog.exec()
        if not dialog.ok:
            return
        for name in dialog.items:
            self.data.del_plot(name)

    def new_item(self):
        itmname = False
        while not itmname:
            itmname, ok = QInputDialog.getText(self, "Enter Name", "Name")
            if not ok:
                return
        itm_energy, ok = QInputDialog.getText(self, "Enter Energy", "Energy")
        if not ok:
            return

        try:
            itm_energy = float(itm_energy)
        except ValueError:
            itm_energy = 0.0

        cname = self.horizontalHeaderItem(self.currentColumn()).text()
        mecobj = self.data.get_mech(cname)
        mecobj.add_itm(itmname, energy=itm_energy)
        self.update_data()

    def chg_unit(self):
        units = ["J", "kJ", "cal", "kcal", "eV", "Ha"]
        dialog = SingleChoiceDialog(units)
        dialog.exec()
        if not dialog.ok:
            return
        unit = dialog.item
        col = self.currentColumn()
        cname = self.horizontalHeaderItem(col).text()
        mobj = self.data.get_mech(cname)
        if mobj is None:
            return
        if unit == mobj.unit:
            return
        mobj.chg_unit(unit)
        self.update_data()

    def contextMenuEvent(self, event):
        menu = QMenu(self)
        freqmenu = menu.addMenu("Frequencies")
        thermomenu = menu.addMenu("Thermochemical Analysis")
        plotmenu = menu.addMenu("Plot")

        addzero = menu.addAction("Add Reference")
        addreac = menu.addAction("Add Reactant")
        settype = menu.addAction("Set Type")
        setpg = menu.addAction("Set Point Group")
        setspin = menu.addAction("Set Spin Multiplicity")
        setmecpprops = menu.addAction("Set MECP Properties")
        editcm = menu.addAction("Edit Comment")
        invertsel = menu.addAction("Invert Selection")
        loadstruct = menu.addAction("Load Structure")
        copytomech = menu.addAction("Copy to Mechanism")
        duplicate = menu.addAction("Duplicate Item")
        mergeitms = menu.addAction("Merge Items")
        update = menu.addAction("Update Item")
        update_data_from_folder = menu.addAction("Update Item From Folder")
        update_selection_from_folder = menu.addAction("Update Selection From Folder")
        export_to_graph = menu.addAction("Export Maplotlib")
        export_csv = menu.addAction("Export Thermochemical CSV")
        removeentropy = menu.addAction("Remove 1/3 of Translational Entropy")
        showrotmat = menu.addAction("Show Rotational Matrix")
        setrotmat = menu.addAction("Set Rotational Matrix")
        importstructvisual = menu.addAction("Import Structure Visual State")

        loadfreqs = freqmenu.addAction("Load Frequencies")
        show_freqs = freqmenu.addAction("Show Frequencies")
        scale_freqs = freqmenu.addAction("Scale Frequencies")
        edit_freqs = freqmenu.addAction("Edit Frequencies")

        calcthermo = thermomenu.addAction("Calculate")
        delthermo = thermomenu.addAction("Delete")
        viewthermoprop = thermomenu.addAction("View Absolute Values")

        newplot = plotmenu.addAction("New Plot")
        addtoplot = plotmenu.addAction("Add to Plot")
        export_plot_excel = plotmenu.addAction("Export Plots to XLSX")

        view_struct_rmol = menu.addAction("Show Structure in MechaEdit")

        action = menu.exec_(self.mapToGlobal(event.pos()))

        if action == addzero:
            self.add_zero()
        elif action == addreac:
            self.add_reac()
        elif action == settype:
            self.set_type()
        elif action == setpg:
            self.set_pg()
        elif action == setspin:
            self.set_spin()
        elif action == setmecpprops:
            self.set_mecp_props()
        elif action == newplot:
            self.new_plot()
        elif action == loadstruct:
            self.load_struct()
        elif action == editcm:
            self.edit_comment()
        elif action == calcthermo:
            self.on_thermo()
        elif action == mergeitms:
            self.on_merge_items()
        elif action == copytomech:
            self.copy_to_mech()
        elif action == export_to_graph:
            self.export_to_matplotlib()
        elif action == loadfreqs:
            self.load_freqs()
        elif action == delthermo:
            self.del_thermo()
        elif action == viewthermoprop:
            self.view_abs_prop()
        elif action == update:
            self.update_itm()
        elif action == update_data_from_folder:
            self.update_itm_data_from_folder()
        elif action == update_selection_from_folder:
            self.update_selection_data_from_folder()
        elif action == invertsel:
            self.invert_selection()
        elif action == export_csv:
            self.export_thermo_csv()
        elif action == show_freqs:
            self.show_freqs()
        elif action == scale_freqs:
            self.scale_freqs()
        elif action == edit_freqs:
            self.edit_freqs()
        elif action == addtoplot:
            self.add_to_plot()
        elif action == duplicate:
            self.duplicate_itm()
        elif action == removeentropy:
            self.remove_entropy()
        elif action == showrotmat:
            self.show_rotmat()
        elif action == setrotmat:
            self.set_rotmat()
        elif action == importstructvisual:
            self.import_struct_visual()
        elif action == export_plot_excel:
            self.export_plot_excel()
        elif action == view_struct_rmol:
            self.view_struct_rmol()

    def view_struct_rmol(self):
        import subprocess as sb
        cellobjs = self.selectedItems()
        colname = self.horizontalHeaderItem(self.currentColumn()).text()
        cmd = shutil.which("rmol")
        if cmd is None:
            cmd = shutil.which("tg")
        if cmd is None:
            diag = QMessageBox(self)
            diag.setText("Could not find mechaedit binary!")
            diag.exec()
            return

        for cellobj in cellobjs:
            itmo = self.data.get_mech(colname).get_itm(cellobj.text())
            with open("temp.xyz", "w") as f:
                f.write(str(itmo.struct.nat) + "\n\n")
                for l, coor in zip(itmo.struct.labels, itmo.struct.coors):
                    f.write(l+"   ")
                    for c in coor:
                        f.write(str(c)+"   ")
                    f.write("\n")

            sb.check_output([cmd, "temp.xyz"])
        sb.call(["rm", "temp.xyz"])

    def export_plot_excel(self):
        plots = self.data.get_plots_names()
        dialog = MultipleChoiceDialog(plots)
        dialog.exec()
        if not dialog.ok:
            return
        plotobjs = self.data.get_plots(dialog.items)

        import openpyxl
        from openpyxl.drawing.text import CharacterProperties, ParagraphProperties, Font
        from openpyxl.chart import (ScatterChart,
                                    Reference,
                                    Series, )

        chart = ScatterChart()
        chart.title = None
        chart.style = 14

        font = Font(typeface='Verdana')
        size = 1200  # 14 point size
        cp = CharacterProperties(latin=font, sz=size, b=False)  # Not bold
        pp = ParagraphProperties(defRPr=cp)
        chart.x_axis.title = 'Reaction Coordinate'
        chart.x_axis.title.tx.rich.p[0].pPr = pp
        chart.y_axis.title = 'Relative Energy'
        chart.y_axis.title.tx.rich.p[0].pPr = pp
        chart.y_axis.majorGridlines = None
        chart.y_axis.minorGridlines = None
        chart.x_axis.majorGridlines = None
        chart.x_axis.minorGridlines = None
        chart.y_axis.scaling.max = round(plotobjs[0].ymax)
        chart.y_axis.scaling.min = round(plotobjs[0].ymin)
        chart.x_axis.scaling.max = round(plotobjs[0].xmax)
        chart.x_axis.scaling.min = round(plotobjs[0].xmin)
        chart.x_axis.tickLblPos = "low"
        chart.legend = None
        stdic = {
            "solid line": "solid",
            "dash line": "dash",
            "dot line": "dot",
            "dash dot line": "dashDot",
            "dash dot dot line": "sysDashDotDot"
        }

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["", ""])
        row = 2
        for plotobj in plotobjs:
            for link in plotobj.links.values():
                # if bar.tp == "link":
                #     continue
                #            _ = ws.cell(column=1, row=row, value="{0:.3f}".format(bar.coors[0]))
                #            _ = ws.cell(column=2, row=row, value="{0:.3f}".format(bar.coors[1]))
                ws.append(link.coors[:2])
                row += 1
                #            _ = ws.cell(column=1, row=row, value="{0:.3f}".format(bar.coors[2]))
                #            _ = ws.cell(column=2, row=row, value="{0:.3f}".format(bar.coors[3]))
                ws.append(link.coors[2:])
                row += 1
                xvalues = Reference(ws, min_col=1, min_row=row - 2, max_row=row - 1)
                values = Reference(ws, min_col=2, min_row=row - 2, max_row=row - 1)
                series = Series(values, xvalues, title_from_data=False)
                color = '%02x%02x%02x' % (link.color[0], link.color[1], link.color[2])
                series.graphicalProperties.line.solidFill = color
                series.graphicalProperties.line.dashStyle = stdic[link.style]
                series.graphicalProperties.line.width = link.width * 13000  # width in EMUs
                # if bar.tp == "bar":
                #     series.graphicalProperties.line.cap = "flat" # width in EMUs
                chart.series.append(series)

            for itm in plotobj.itms.values():
                ws.append(itm.coors[:2])
                row += 1
                ws.append(itm.coors[2:])
                row += 1
                xvalues = Reference(ws, min_col=1, min_row=row - 2, max_row=row - 1)
                values = Reference(ws, min_col=2, min_row=row - 2, max_row=row - 1)
                series = Series(values, xvalues, title_from_data=False)
                color = '%02x%02x%02x' % (itm.color[0], itm.color[1], itm.color[2])
                series.graphicalProperties.line.solidFill = color
                series.graphicalProperties.line.dashStyle = stdic[itm.style]
                series.graphicalProperties.line.width = itm.width * 13000  # width in EMUs
                chart.series.append(series)
        ws.add_chart(chart, "C2")
        wb.save("plots.xlsx")

    def import_struct_visual(self):
        cellobjs = self.selectedItems()
        colname = self.horizontalHeaderItem(self.currentColumn()).text()
        diag = SingleChoiceDialog(self.data.get_mech(colname).get_itms_names())
        diag.exec()
        if not diag.item:
            return
        itmo1 = self.data.get_mech(colname).get_itm(diag.item)

        for cellobj in cellobjs:
            itmo = self.data.get_mech(colname).get_itm(cellobj.text())
            if itmo is None:
                continue
            for i,_ in enumerate(itmo.struct.rot_matrix):
                for j in range(3):
                    itmo.struct.rot_matrix[i][j] = itmo1.struct.rot_matrix[i][j]
            itmo.struct.rRange = itmo1.struct.rRange
            itmo.struct.lRange = itmo1.struct.lRange
            itmo.struct.dRange = itmo1.struct.dRange
            itmo.struct.uRange = itmo1.struct.uRange

    def set_rotmat(self):
        rotmat, ok = QInputDialog.getMultiLineText(self, "Enter rotation matrix", "Enter rotation matrix")
        if not ok:
            return
        rotmat = rotmat.splitlines()
        for i in range(3):
            rotmat[i] = [float(j) for j in rotmat[i].split()]

        colname = self.horizontalHeaderItem(self.currentColumn()).text()
        cellobjs = self.selectedItems()
        for cellobj in cellobjs:
            itmname = cellobj.text()
            itmo = self.data.get_mech(colname).get_itm(itmname)
            if itmo is None:
                continue
            itmo.struct.rot_matrix = rotmat

    def show_rotmat(self):
        colname = self.horizontalHeaderItem(self.currentColumn()).text()
        cellobj = self.selectedItems()[0]
        itmname = cellobj.text()
        itmo = self.data.get_mech(colname).get_itm(itmname)
        if itmo is None:
            return
        rotmat = ""
        for i in itmo.struct.rot_matrix:
            for j in i:
                rotmat += str(j) + "   "
            rotmat += " \n"
        diag = QMessageBox(self)
        diag.setText(rotmat)
        diag.exec()

    def remove_entropy(self):
        for cobj in self.selectedItems():
            rowmane = self.verticalHeaderItem(cobj.row()).text()
            if rowmane != "Name":
                continue
            cname = self.horizontalHeaderItem(cobj.column()).text()
            itmobj = self.data.get_mech(cname).get_itm(cobj.text())
            if itmobj is None:
                continue
            for t in itmobj.temps:
                print(itmobj.thermo["strans"][t])
                itmobj.thermo["strans"][t] *= 2/3
                print(itmobj.thermo["strans"][t])

    def duplicate_itm(self):
        cobj = self.selectedItems()[0]
        cname = self.horizontalHeaderItem(cobj.column()).text()
        newitmobj = self.data.get_mech(cname).duplicate_itm(cobj.text())
        if newitmobj is None:
            msg = QMessageBox()
            msg.setText("Item could not be duplicated, it does not exist!")
            msg.exec()
            return
        self.update_data()

    def add_to_plot(self):
        col = self.currentColumn()
        colname = self.horizontalHeaderItem(col).text()
        mec = self.data.get_mech(colname)

        # cojo los objectos celda seleccinados e itero sobre ellos
        # tomando el texto de la celda
        # con ese texto obtengo los objectos intermedios que devolveran None
        # si no existe ninguno con el nombre del texto de la celda
        cellobjs = self.selectedItems()
        itmes = [mec.get_itm(cb.text()) for cb in cellobjs if mec.get_itm(cb.text()) is not None]
        itmes = [itm for itm in itmes if itm is not None or itm.tp != "ref"]
        commonrefs = mec.itms_common_refs(itmes)
        # commonrefs = [0]
        if not commonrefs:
            msg = QMessageBox()
            msg.setText("No common reference for selected intermediates to set up a plot")
            msg.exec()
            return
        elif len(commonrefs) == 1:
            plots = self.data.get_plots_names()
            dialog = SingleChoiceDialog(plots)
            dialog.exec()
            if not dialog.ok:
                return
            plotobj = self.data.get_plot(dialog.item)
            if plotobj is None:
                return
            if plotobj.ref != commonrefs[0]:
                msg = QMessageBox()
                msg.setText("Selected plot and items have different references")
                msg.exec()
                return

            plotobj.add_itms(itmes)
            # self.data.add_plot(gname, commonrefs[0], itmes)
            self.parent.plot.update_data()

    def edit_freqs(self):
        cobj = self.selectedItems()[0]
        cname = self.horizontalHeaderItem(cobj.column()).text()
        itmo = self.data.get_mech(cname).get_itm(cobj.text())
        if itmo is None:
            return
        txt = " \n".join([str(f) for f in itmo.freqs.vibs])
        txt += "\n"+str(itmo.freqs.unit)
        dialog, ok = QInputDialog.getMultiLineText(self,
                                                   "Edit Frequencies",
                                                   "Last line is the unit", txt)
        if not ok:
            return

        freqs = dialog.splitlines()
        unit =freqs[-1]
        if unit not in ["cm-1"]:
            print("bad unit for frequencies!")
            return
        vibs = [float(f) for f in freqs[:-1]]
        itmo.freqs.vibs = vibs
        itmo.freqs.unit = unit

    def scale_freqs(self):
        coeff = False
        while not coeff:
            coeff, ok = QInputDialog.getText(self, "Scale Frequencies", "Enter Scale Factor")
            if not ok:
                return
            try:
                coeff = float(coeff)
            except ValueError:
                coeff = False

        for cellobj in self.selectedItems():
            rname = self.verticalHeaderItem(cellobj.row()).text()
            if rname != "Name":
                continue
            cname = self.horizontalHeaderItem(cellobj.column()).text()
            itmname = cellobj.text()
            itmobj = self.data.get_mech(cname).get_itm(itmname)
            itmobj.scale_freqs(coeff)

    def show_freqs(self):
        cobj = self.selectedItems()[0]
        print("cobj", cobj.text())  
        cname = self.horizontalHeaderItem(cobj.column()).text()
        itmo = self.data.get_mech(cname).get_itm(cobj.text())
        print("itmo", self.data.get_mech(cname))
        if itmo is None:
            return
        dialog = QMessageBox()
        print("itmo.freqs.vibs", itmo.freqs.vibs)
        if itmo.merged is True:
            txt = ""
            for it in itmo.itms:
                print("it", it.name)
                txt += " \n".join([str(f) for f in it.freqs.vibs])
                txt += "\n"
                print("txt", txt)
        else:
            txt = " \n".join([str(f) for f in itmo.freqs.vibs])
            txt += "\n"+itmo.freqs.unit

        
        dialog.setText(txt)
        dialog.exec()

    def export_thermo_csv(self):
        fname, ext = QFileDialog.getSaveFileName()
        if not fname:
            return
        f = open(fname + ".csv", "w")
        cellobjs = self.selectedItems()
        for co in cellobjs:
            rname = self.verticalHeaderItem(co.row()).text()
            if rname != "Name":
                continue
            cname = self.horizontalHeaderItem(co.column()).text()
            itmo = self.data.get_mech(cname).get_itm(co.text())
            if itmo is None:
                continue
            temps = itmo.temps
            reacobjs = itmo.get_reacs()
            if not reacobjs:
                continue
            reacprops = reacobjs[0].thermo.keys()
            f.write('"'+itmo.name+'"' + ";\n")
            for reacobj in reacobjs:
                f.write('"'+reacobj.name+'"' + ";")
                f.write(";".join([r.upper() for r in reacprops]) + ";\n")
                for temp in temps:
                    if temp not in reacobj.temps:
                        continue
                    f.write(str(temp) + ";")
                    for rprop in reacprops:
                        strval = "{:.4e}".format(reacobj.thermo[rprop][temp])
                        strval += ";"
                        f.write(strval)
                    f.write(";\n")
            f.write(";\n")
        f.close()

    def invert_selection(self):
        cellobjs = self.selectedItems()
        col = self.currentColumn()
        for row in range(self.rowCount()):
            cellobj = self.item(row, col)
            if cellobj in cellobjs:
                cellobj.setSelected(False)
            else:
                if self.verticalHeaderItem(row).text() == "Name":
                    cellobj.setSelected(True)

    def update_itm(self):
        cellobjs = self.selectedItems()
        for cellobj in cellobjs:
            try:
                col, row = cellobj.column(), cellobj.row()
            except RuntimeError:
                continue
            rowname = self.verticalHeaderItem(row).text()

            if rowname != "Name":
                continue
            colname = self.horizontalHeaderItem(col).text()
            itmname = cellobj.text()
            itmobj = self.data.get_mech(colname).get_itm(itmname)
            if itmobj is None:
                continue
            itmobj.full_update()
            self.update_data()

    def update_itm_data_from_folder(self):
        """
        This is design for only one item. There it will take the first one selected
        """
        cellobjs = self.selectedItems()
        if not cellobjs:
            Log.show_msg("You must select one item (in a row with the label 'Name')")
            return 
        cellobj = cellobjs[0]

        try:
            col, row = cellobj.column(), cellobj.row()
        except RuntimeError:
            return
        rowname = self.verticalHeaderItem(row).text()

        if rowname != "Name":
            Log.show_msg("You must select the name of the item")
            return


        colname = self.horizontalHeaderItem(col).text()
        itmname = cellobj.text()
        itmobj = self.data.get_mech(colname).get_itm(itmname)
        if itmobj is None:
            return

        # get file name
        dialog = QFileDialog()
        dialog.setFileMode(QFileDialog.Directory)
        dialog.setOption(QFileDialog.DontUseNativeDialog)
        folder = dialog.getExistingDirectory(self, "Choose folder",
                                                  self.parent.last_dir)
        print(os.path.basename(str(folder)))
        if os.path.basename(str(folder)) != itmname:
            Log.show_msg("To avoid mistakes, the name of the folder must equal the name of the item")
            return
        
        # aqui tengo que poner el dialog para coger el directory
        mecobj = self.data.get_mech(colname)
        itm = mecobj.itm_from_folder(folder, add=False)
        if itm is None:
            Log.show_msg("Could not load data from folder: "+str(itm))
            return
        itmobj.update_data_from_itm(itm)
        self.update_data()

    def update_selection_data_from_folder(self):
        """
        This is design for only one item. There it will take the first one selected
        """
        cellobjs = self.selectedItems()
        if not cellobjs:
            Log.show_msg("You must select one item (in a row with the label 'Name')")
            return 

        # get file name
        dialog = QFileDialog()
        dialog.setFileMode(QFileDialog.Directory)
        dialog.setOption(QFileDialog.DontUseNativeDialog)
        folder = dialog.getExistingDirectory(self, "Choose folder",
                                                  self.parent.last_dir)
        if not str(folder):
            return
        try:
            dirs = os.listdir(str(folder))
        except:
            Log.show_msg("File not found")
            return

        for cellobj in cellobjs:
            try:
                col, row = cellobj.column(), cellobj.row()
            except RuntimeError:
                return
            rowname = self.verticalHeaderItem(row).text()

            if rowname != "Name":
                #Log.show_msg("You must select the name of the item")
                continue

            colname = self.horizontalHeaderItem(col).text()
            itmname = cellobj.text()
            itmobj = self.data.get_mech(colname).get_itm(itmname)
            if itmobj is None:
                continue

            # find folder
            itmdir = None
            for d in dirs:
                if d == itmname:
                    itmdir = d
                    break
            if itmdir is None:
                continue
            
            # aqui tengo que poner el dialog para coger el directory
            mecobj = self.data.get_mech(colname)
            itm = mecobj.itm_from_folder(os.path.join(str(folder), itmdir), add=False)
            if itm is None:
                Log.show_msg("Could not load data from folder: "+str(itm))
                continue
            itmobj.update_data_from_itm(itm)
            print("Updated item: ", itmname)
        self.update_data()

    def on_header_merge_items(self):
        col = self.currentColumn()
        colname = self.horizontalHeaderItem(col).text()
        mecobj = self.data.get_mech(colname)
        #itmlist = mecobj.get_itms_names(tp="min")
        itmlist = mecobj.get_itms_names()
        dialog = MergeItemsDialog(itmlist)
        dialog.exec()
        if dialog.ok:
            try:
                if dialog.fullmerge:
                    mecobj.fullmerge_itms(dialog.items)
                else:
                    mecobj.merge_itms(dialog.items)
            except TypeError as e:
                msg = QMessageBox()
                msg.setText("Could not merge items \n" + str(e))
                msg.exec()
            except ValueError as e:
                msg = QMessageBox()
                msg.setText("Could not merge items \n" + str(e))
                msg.exec()
            else:
                self.update_data()

    def view_abs_prop(self):
        cellobj = self.selectedItems()[0]
        rname = self.verticalHeaderItem(cellobj.row()).text()
        if rname != "Name":
            return
        cname = self.horizontalHeaderItem(cellobj.column()).text()
        itmname = cellobj.text()
        itmobj = self.data.get_mech(cname).get_itm(itmname)
        if itmobj is None:
            return
        props = itmobj.thermo.keys()
        msg = "    ".join(props)
        msg += "  \n"
        for temp in itmobj.temps:
            msg += str(temp) + ":  "
            for prop in props:
                msg += "{:.3f}".format(itmobj.thermo[prop][temp]) + "   "
            msg += "\n"
        msgbox = QMessageBox(self)
        msgbox.setText(msg)
        msgbox.exec()

    def del_thermo(self):
        for cellobj in self.selectedItems():
            colname = self.horizontalHeaderItem(cellobj.column()).text()
            rowname = self.verticalHeaderItem(cellobj.row()).text()
            if rowname != "Name":
                continue
            itmobj = self.data.get_mech(colname).get_itm(cellobj.text())
            if itmobj is None:
                continue
            itmobj.del_thermo()

    def load_freqs(self):
        cellobj = self.selectedItems()[0]
        colname = self.horizontalHeaderItem(cellobj.column()).text()
        itmobj = self.data.get_mech(colname).get_itm(cellobj.text())
        if itmobj is None:
            return
        f, _ = QFileDialog.getOpenFileName()
        if f:
            if not itmobj.freqs_from_file(f):
                msb = QMessageBox()
                msb.setText("Could not load frequencies")
                msb.exec()

    def export_to_matplotlib(self):
        cellobjs = self.selectedItems()
        itmnames = [cellobj.text() for cellobj in cellobjs
                    if self.verticalHeaderItem(cellobj.row()).text() == "Name"]

        f = open("matinput", "a")
        colname = self.horizontalHeaderItem(cellobjs[0].column()).text()
        pos = 0
        f.write("0 zero%       0      blue_b\n")
        for num, itmname in enumerate(itmnames):
            itmobj = self.data.get_mech(colname).get_itm(itmname)
            if itmobj is None:
                continue
            pos += 1
            f.write(str(pos) + "     ")
            f.write(itmobj.name + "     ")
            f.write(str(itmobj.get_def_ref_energy()) + "    ")
            if num < len(itmnames) - 1:
                f.write("blue_b\n")
            else:
                f.write("n\n")
        f.close()

    def copy_to_mech(self):
        allmechs = self.data.get_mechs_names()
        dialog = MultipleChoiceDialog(allmechs)
        dialog.exec()
        if not dialog.ok:
            return
        cellobjs = self.selectedItems()
        for cellobj in cellobjs:
            col, row = cellobj.column(), cellobj.row()
            colname = self.horizontalHeaderItem(col).text()
            rowname = self.verticalHeaderItem(row).text()
            if rowname != "Name":
                continue
            mecobj = self.data.get_mech(colname)
            itmobj = mecobj.get_itm(cellobj.text())
            if itmobj is None:
                continue
            for mecname in dialog.items:
                dmecobj = self.data.get_mech(mecname)
                dmecobj.add_itm(itmobj)
        self.update_data()

    def on_merge_items(self):
        cellobjs = self.selectedItems()
        cols = [cellobj.column() for cellobj in cellobjs]

        # vertify that there are cells from only one column
        if len(set(cols)) > 1:
            msg = QMessageBox()
            msg.setText("Please select items from only one column")
            msg.exec()
            return

        colname = self.horizontalHeaderItem(cols[0]).text()
        mecobj = self.data.get_mech(colname)
        itmlist = [cellobj.text() for cellobj in cellobjs
                   if cellobj is not None and
                   self.verticalHeaderItem(cellobj.row()).text() == "Name"]
        try:
            mecobj.merge_itms(itmlist)
        except TypeError as e:
            msg = QMessageBox()
            msg.setText("Could not merge items \n" + str(e))
            msg.exec()
        except ValueError as e:
            msg = QMessageBox()
            msg.setText("Could not merge items \n" + str(e))
            msg.exec()
        else:
            self.update_data()

    def on_thermo(self):
        ok, temps, modes, P, V, nnfreq = ThermoDialog(self).get_values()

        if not ok:
            return
        cellobjs = self.selectedItems()
        for cellobj in cellobjs:
            col, row = cellobj.column(), cellobj.row()
            colname = self.horizontalHeaderItem(col).text()
            rowname = self.verticalHeaderItem(row).text()
            if rowname == "Name":
                itmname = cellobj.text()
                itmobj = self.data.get_mech(colname).get_itm(itmname)
                if itmobj is None:
                    continue
                #     verifying if the merged items are OK---------------------
                if itmobj.merged:
                    # check for freqs in items
                    itmnoref = [i.name for i in itmobj.get_itms()
                                if not i.freqs.vibs]
                    if itmnoref:
                        msg = QMessageBox()
                        msg.setText("No frequencies for items:" + ", ".join(itmnoref) + "\nDo you want to continue?")
                        msg.setStandardButtons(QMessageBox.Yes)
                        msg.addButton(QMessageBox.No)
                        if msg.exec() == QMessageBox.No:
                            continue
                    # check for point group in items
                    next_itm = False
                    for i in itmobj.get_itms():
                        if i.pg is None:
                            ok, pg = PGDialog(self, i.name).get_values()
                            if not ok:
                                next_itm = True
                                break
                            itmobj.pg = pg
                    if next_itm:
                        continue
                    # check compatibility between freqs, struct and point group
                    next_itm = False
                    for i in itmobj.get_itms():
                        if i.pg != "solid" and not i.freqs.vibs or \
                        i.pg != "solid" and not i.struct.atoms:
                            msg = QMessageBox()
                            text = "This Point group:"
                            text += itmobj.pg
                            text += " of item: "
                            text += itmobj.name
                            text += " requires frequencies and coordinates"
                            msg.setText(text)
                            msg.exec()
                            next_itm = True
                            break
                    if next_itm:
                        continue
                #     verifying if the merged items are OK---------------------
                else:
                    # check for freqs
                    if not itmobj.freqs.vibs:
                        msg = QMessageBox()
                        msg.setText("No frequencies for " + itmobj.name + "\nDo you want to continue?")
                        msg.setStandardButtons(QMessageBox.Yes)
                        msg.addButton(QMessageBox.No)
                        if msg.exec() == QMessageBox.No:
                            continue
                    # verify pg
                    if itmobj.pg is None:
                        ok, pg = PGDialog(self, itmobj.name).get_values()
                        if not ok:
                            continue
                        itmobj.pg = pg

                    if itmobj.pg != "solid" and not itmobj.freqs.vibs or \
                    itmobj.pg != "solid" and not itmobj.struct.atoms:
                        msg = QMessageBox()
                        text = "This Point group:"
                        text += itmobj.pg
                        text += " of item: "
                        text += itmobj.name
                        text += " requires frequencies and coordinates"
                        msg.setText(text)
                        msg.exec()
                        continue
                try:
                    # If single MECP item (not merged), prompt freq selection in GUI
                    restored = False

                    itmobj.calc_thermo(temps, modes, P, V, nnfreq)
                    # restore original freqs if we changed them
                    if restored and hasattr(itmobj, "_orig_freqs"):
                        itmobj.freqs.vibs = itmobj._orig_freqs
                        del itmobj._orig_freqs
                except Exception as e:
                    print(e)
                    msg = QMessageBox()
                    errorm = "There was a problem computing: "
                    errorm += itmobj.name
                    errorm += " " + str(e)
                    errorm += "\n Please check the point group"
                    msg.setText(errorm)
                    msg.exec()
            else:
                continue

    def edit_comment(self, cm=""):
        cellobj = self.selectedItems()[0]
        col, row = cellobj.column(), cellobj.row()
        colname = self.horizontalHeaderItem(col).text()
        rowname = self.verticalHeaderItem(row).text()
        if rowname != "Name":
            return
        itmname = cellobj.text()
        itmobj = self.data.get_mech(colname).get_itm(itmname)
        if itmobj is None:
            print("Bad intermediate")
            return
        if cm:
            itmobj.cm = cm
        else:
            cm, ok = QInputDialog.getMultiLineText(self, "Edit Comment",
                                                   "Comment", itmobj.cm)
            if ok:
                itmobj.cm = cm

    def load_struct(self):
        infile = QFileDialog.getOpenFileName()
        if not infile[0]:
            return
        cellobj = self.selectedItems()[0]
        row, col = cellobj.row(), cellobj.column()
        rowname = self.verticalHeaderItem(row).text()
        if rowname == "Name":
            itmname = cellobj.text()
            colname = self.horizontalHeaderItem(col).text()
            itmobj = self.data.get_mech(colname).get_itm(itmname)
            itmobj.struct_from_file(infile[0])

    def new_plot(self):
        col = self.currentColumn()
        colname = self.horizontalHeaderItem(col).text()
        mec = self.data.get_mech(colname)

        # cojo los objectos celda seleccinados e itero sobre ellos
        # tomando el texto de la celda
        # con ese texto obtengo los objectos intermedios que devolveran None
        # si no existe ninguno con el nombre del texto de la celda
        cellobjs = self.selectedItems()
        itm_names = [cb.text() for cb in cellobjs if mec.get_itm(cb.text()) is not None]
        # get items objects
        itms = [mec.get_itm(itm_name) for itm_name in itm_names if mec.get_itm(itm_name) is not None]

        # filter out refs and None
        itms = [itm for itm in itms if itm is not None or itm.tp != "ref"]

        commonrefs = mec.itms_common_refs(itms)
        if not commonrefs:
            msg = QMessageBox()
            msg.setText("No common reference for selected intermediates to set up a plot")
            msg.exec()
            return
        elif len(commonrefs) == 1:
            plotname, ok = QInputDialog().getText(self, "Enter New Name: ", "Plot Name:")
            if not ok:
                return
            self.data.add_plot(plotname, commonrefs[0], colname, itms, labels=True)
            self.parent.plot.update_data()

    def add_ref_to_mech(self):
        col = self.currentColumn()
        colname = self.horizontalHeaderItem(col).text()
        mec = self.data.get_mech(colname)
        zerolist = mec.get_itms_names(tp="ref")
        dialog = MultipleChoiceDialog(zerolist)
        dialog.exec()
        if not dialog.ok:
            return
        refs = []
        for ref in dialog.items:
            ref = mec.get_itm(ref)
            if ref is None or len(ref.struct.atoms) == 0:
                msg = QMessageBox()
                msg.setText(ref.name + " does not have atoms! Please set atoms\n")
                msg.exec()
                return
            refs.append(ref)
        refname, ok = QInputDialog.getText(self, "Name", "Enter Name")
        if not ok:
            return
        try:
            mec.add_ref(refname, refs)
        except NameError as e:
            msg = QMessageBox()
            msg.setText(ref.name + " Already exists! \n")
            msg.exec()
            return

    def set_def_basis_ref(self):
        col = self.currentColumn()
        colname = self.horizontalHeaderItem(col).text()
        mec = self.data.get_mech(colname)
        zerolist = mec.get_itms_names(tp="ref")
        dialog = MultipleChoiceDialog(zerolist)
        dialog.exec()
        if dialog.ok:
            refs = []
            for ref in dialog.items:
                ref = mec.get_itm(ref)
                if len(ref.struct.atoms) == 0:
                    msg = QMessageBox()
                    msg.setText(ref.name + " does not have atoms! Please set atoms\n")
                    msg.exec()
                    return
                refs.append(ref)
            mec.set_def_basis_ref(refs)
            # for itmobj in mec.get_itms():
            #     if itmobj.tp != "ref":
            #         if len(itmobj.atoms) > 1:
            #             itmobj.add_ref(None, refs)

    def get_current_itm_name(self):
        itmobj = self.currentItem()
        return itmobj.text()

    def del_itms(self):
        for cellobj in self.selectedItems():
            row, col = cellobj.row(), cellobj.column()
            rowname = self.verticalHeaderItem(row).text()
            if rowname != "Name":
                continue
            itmname = cellobj.text()
            colname = self.horizontalHeaderItem(col).text()
            itmobj = self.data.get_mech(colname).get_itm(itmname)
            if itmobj.tp == "ref":
                msg = QMessageBox()
                msg.setText("Deleting reference is not implemented")
                msg.exec()
                return
                ref_names = ""
                for ref in self.data.get_mech(colname).get_refs():
                    for itm in ref.itms:
                        if itm is itmobj:
                            ref_names += ref.name + "\n"
                if ref_names:
                    msg = QMessageBox()
                    msg.setText("This item cannot be deleted without deleting:\n" + ref_names)
                    msg.exec()
            else:
                for itm in self.data.get_mech(colname).get_itms():
                    if itm.tp != "ref" and itm is not itmobj:
                        for reac in itm.get_reacs():
                            if reac.ref is itmobj:
                                print(reac.name, " will be deleted from ", itm.name)
                                itm.del_reac(reac)

            self.data.get_mech(colname).del_itm(itmname)
        # self.editing = False
        self.update_data()

    def set_type(self):
        refsnames = ["min", "ts", "ref", "mecp"]
        dialog = MultipleChoiceDialog(refsnames)
        dialog.exec()
        

        cellobjs = self.selectedItems()
        for cellobj in cellobjs:
            col, row = cellobj.column(), cellobj.row()
            rowname = self.verticalHeaderItem(row).text()
            if rowname != "Name":
                continue
            colname = self.horizontalHeaderItem(col).text()
            itmname = cellobj.text()
            itmobj = self.data.get_mech(colname).get_itm(itmname)
            if dialog.ok:
                if dialog.items:
                    tp = dialog.items[0]
                    itmobj.tp = tp
                else:
                    return
            else:
                return
            
            itmobj.update_rel_thermo()

    def set_pg(self):
        pgs = ["solid"] + list(Thermo.sigma.keys())
        dialog = SingleChoiceDialog(pgs)
        dialog.exec()
        if dialog.ok:
            pg = dialog.item
        else:
            return
        cellobjs = self.selectedItems()
        for cellobj in cellobjs:
            col, row = cellobj.column(), cellobj.row()
            rowname = self.verticalHeaderItem(row).text()
            if rowname != "Name":
                continue
            colname = self.horizontalHeaderItem(col).text()
            itmname = cellobj.text()
            itmobj = self.data.get_mech(colname).get_itm(itmname)
            if itmobj is None:
                return
            itmobj.pg = pg
            # ver si tengo que recalcular thermo con el pg, deberia
            # itmobj.update_rel_thermo()

    def set_spin(self):
        spin = False
        while not spin:
            spin, ok = QInputDialog.getText(self, "Set spin", "Set spin")
            if not ok:
                return
            try:
                spin = float(spin)
            except ValueError:
                spin = False

        cellobjs = self.selectedItems()
        for cellobj in cellobjs:
            col, row = cellobj.column(), cellobj.row()
            rowname = self.verticalHeaderItem(row).text()
            if rowname != "Name":
                continue
            colname = self.horizontalHeaderItem(col).text()
            itmname = cellobj.text()
            itmobj = self.data.get_mech(colname).get_itm(itmname)
            if itmobj is None:
                return
            itmobj.spin = spin

    def set_mecp_props(self):
        # gather defaults from the first selected MECP "Name" row (if any)
        red_def = SOC_def = grad_def = diffgrad_def = None
        sel = self.selectedItems()
        for cellobj in sel:
            try:
                col, row = cellobj.column(), cellobj.row()
            except RuntimeError:
                continue
            vh = self.verticalHeaderItem(row)
            if vh is None or vh.text() != "Name":
                continue
            colname = self.horizontalHeaderItem(col).text()
            itmname = cellobj.text()
            itmobj = self.data.get_mech(colname).get_itm(itmname)
            if itmobj is None:
                continue
            if getattr(itmobj, "tp", None) != "mecp":
                continue
            # use existing values as defaults when present (not None)
            if getattr(itmobj, "redmass", None) is not None:
                red_def = itmobj.redmass
            if getattr(itmobj, "SOC", None) is not None:
                SOC_def = itmobj.SOC
            else:
                SOC_def = 100.0  # a common default
            if getattr(itmobj, "grad", None) is not None:
                grad_def = itmobj.grad
            if getattr(itmobj, "diffgrad", None) is not None:
                diffgrad_def = itmobj.diffgrad
            break

        dialog = MECPropsDialog(self, redmass=red_def, SOC=SOC_def, grad=grad_def, diffgrad=diffgrad_def)
        ok, redmass, SOC, grad, diffgrad = dialog.get_values()
        if not ok:
            return

        # apply values to all selected Name rows (validate each is a mecp)
        for cellobj in sel:
            try:
                col, row = cellobj.column(), cellobj.row()
            except RuntimeError:
                continue
            rowname = self.verticalHeaderItem(row).text()
            if rowname != "Name":
                continue
            colname = self.horizontalHeaderItem(col).text()
            itmname = cellobj.text()
            itmobj = self.data.get_mech(colname).get_itm(itmname)
            if itmobj is None:
                return
            if itmobj.tp != "mecp":
                msg = QMessageBox()
                msg.setText(itmobj.name + " is not a mecp!")
                msg.exec()
                return
            itmobj.redmass = redmass
            itmobj.SOC = SOC
            itmobj.grad = grad
            itmobj.diffgrad = diffgrad

    def import_itm(self):
        dialog = QFileDialog()
        dialog.setFileMode(QFileDialog.Directory)
        dialog.setOption(QFileDialog.DontUseNativeDialog)
        folders = dialog.getExistingDirectory(self, "Choose folder",
                                                  self.parent.last_dir)
        folders = [folders]
        # file_dialog = QFileDialog()
        # file_dialog.setFileMode(QFileDialog.DirectoryOnly)
        # file_dialog.setOption(QFileDialog.DontUseNativeDialog, True)
        # file_view = file_dialog.findChild(QListView, 'listView')

        # to make it possible to select multiple directories:
        # if file_view:
        #     file_view.setSelectionMode(QAbstractItemView.MultiSelection)
        # f_tree_view = file_dialog.findChild(QTreeView)
        # if f_tree_view:
        #     f_tree_view.setSelectionMode(QAbstractItemView.MultiSelection)
        #
        # if not file_dialog.exec():
        #     return
        # folders = file_dialog.selectedFiles()
        ignored_itms = []
        col = self.currentColumn()
        colname = self.horizontalHeaderItem(col).text()
        self.setCursor(QCursor(Qt.WaitCursor))
        for folder in folders:
            itmobj = self.data.get_mech(colname).itm_from_folder(folder)
            # try:
            #     itmobj = self.data.get_mech(colname).itm_from_folder(folder)
            # except Exception as e:
            #     print(e)
            #     continue
            if itmobj is None:
                ignored_itms.append(folder)
        self.unsetCursor()
        self.parent.last_dir = os.path.basename(folders[0])
        self.update_data()
        self.parent.last_dir = folders[0]

        # inform bad itms
        if ignored_itms:
            msg = "The following folders were not added because there are items with the same name or folder is not set properly:\n\n"
            msg += "\n".join([i for i in ignored_itms])
            Log.show_msg(msg)


    def add_zero(self):
        set_def_ref = False
        # dictionary of references from mech
        # every element of the dic is a list of Itm objects with which the ref is made up
        refs_from_mec = {}
        newreflist = []  # list of ref itms to make a new reference
        newrefname = ""

        cellobjs = self.selectedItems()
        cols = [cellobj.column() for cellobj in cellobjs]

        # vertify that there are cells from only one column
        if len(set(cols)) > 1:
            Log.show_msg("Please select items from only one column")
            return

        # get the label of the column using the first index of the list
        # cols since they all have to be the same
        colname = self.horizontalHeaderItem(cols[0]).text()
        mecobj = self.data.get_mech(colname)

        # show dialog with the ref itms and references of the mechanism
        #  and the default referece
        refsnames = mecobj.get_itms_names(tp="ref")
        refsnames += list(mecobj.refs.keys())
        refsnames += ["Default Reference"]
        dialog = MultipleChoiceDialog(refsnames)
        dialog.exec()
        if not dialog.ok:
            return
        # --------------------------------------------------------------

        # classify selection to see what to do refs list
        for ref in dialog.items:
            if ref == "Default Reference":
                set_def_ref = True
            elif ref in mecobj.refs.keys():
                refs_from_mec[ref] = mecobj.refs[ref]
            else:
                newreflist.append(ref)
        # ---------------------------------------------------------------

        # -----------------Build the ref with newreflist ------------------------------------------
        refs = []
        if newreflist:
            newrefname, ok = QInputDialog.getText(self, "Enter Name", "New Reference Name")
            if ok and newrefname:
                #  build and add new ref from list of unclassified items-----------------------------
                for ref in newreflist:
                    ref = self.data.get_mech(colname).get_itm(ref)
                    if len(ref.struct.atoms) == 0:
                        msg = QMessageBox()
                        msg.setText(ref.name + " does not have atoms! Please set atoms\n")
                        msg.exec()
                        return
                    refs.append(ref)
                # try to add the new ref to mechanism mecobj
                try:
                    if mecobj.add_ref(newrefname, refs) is None:
                        msg = QMessageBox()
                        msg.setText(newrefname + " Not added to mechanims\n Will not be added to items either\n")
                        msg.exec()
                except NameError:
                    pass
                    #     try to add the new ref to mechanism mecobj
        # -----------------Build the ref with newreflist ------------------------------------------

        for cellobj in cellobjs:
            if cellobj is None:
                continue
            row = cellobj.row()
            rowname = self.verticalHeaderItem(row).text()
            if rowname != "Name":
                continue

            itmname = cellobj.text()
            itmobj = mecobj.get_itm(itmname)
            if itmobj is None or itmobj.tp == "ref":
                continue

            if len(itmobj.struct.atoms) == 0:
                msg = QMessageBox()
                msg.setText(itmobj.name + " does not have atoms! Please set atoms\n")
                msg.exec()
                continue

            # start adding references in the lists built from selection
            if set_def_ref:
                itmobj.add_default_reference()
            for refname, value in refs_from_mec.items():
                try:
                    ref = itmobj.add_ref(refname, value)
                    if ref is None:
                        msg = QMessageBox()
                        msg.setText("Item " + itmobj.name + "already has a reference with that name!")
                        msg.exec()
                # except np.linalg.LinAlgError as e:  apparently there is also index error
                except Exception as e:
                    msg = QMessageBox()
                    msg.setText("Error: "+str(e)+"\nDo you want to enter the coefficients manually?")
                    msg.setStandardButtons(QMessageBox.Ok | QMessageBox.Cancel)
                    ret = msg.exec()
                    if ret == QMessageBox.Ok:
                        coefstring = "   ".join([it.name for it in value])
                        coefs, ok = QInputDialog.getText(self, "Enter coefficients for "+itmobj.name,
                                                         itmobj.name+"\nCoeficients : "+coefstring)
                        if ok and coefs:
                            coefs = coefs.split()
                            coefs = [float(c) for c in coefs]
                            itmobj.add_ref(refname, value, coefs)

            # Finally add the new ref to the itmobj i
            if refs:
                try:
                    if itmobj.add_ref(newrefname, refs) is None:
                        msg = QMessageBox()
                        msg.setText("Item " + itmobj.name + "already has a reference with that name!")
                        msg.exec()
                        continue
                # except ValueError as e: apparently there is also index error
                except Exception as e:
                    msg = QMessageBox()
                    msg.setText(str(e))
                    msg.exec()
                    continue

    def add_reac(self):
        currentColumn = self.currentColumn()
        colname = self.horizontalHeaderItem(currentColumn).text()
        itmsnames = self.data.get_mech(colname).get_itms_names()
        dialog = MultipleChoiceDialog(itmsnames)
        dialog.exec()
        if not dialog.ok:
            return
        for cellobj in self.selectedItems():
            col, row = cellobj.column(), cellobj.row()
            if col != currentColumn:
                continue
            rowname = self.verticalHeaderItem(row).text()
            if rowname != "Name":
                continue
            itmname = cellobj.text()
            itmobj = self.data.get_mech(colname).get_itm(itmname)
            for reac in dialog.items:
                reacobj = self.data.get_mech(colname).get_itm(reac)
                if reacobj.struct.nat != itmobj.struct.nat:
                    if len(reacobj.struct.atoms) != len(itmobj.struct.atoms):
                        msg = QMessageBox()
                        text = "Reactant reference and itermediate do not have" \
                               " the same numer of types of atoms"
                        msg.setText(text)
                        msg.exec()
                        return

                    # print("atoms len different")
                    # check if they have the same references to calculate reaction values
                    # based on that reference
                    same_refs_names = []
                    for itmref in itmobj.get_refs():
                        for reacref in reacobj.get_refs():
                            if itmref.name == reacref.name:
                                same_refs_names.append(itmref.name)
                    if not same_refs_names:
                        msg = QMessageBox()
                        text = "Reactant reference and itermediate do not have" \
                               " the same number of atoms, neither coincident" \
                               " references. \nPlease add same reference to both"
                        msg.setText(text)
                        msg.exec()
                        return
                    elif len(same_refs_names) == 1:
                        print("len 11")
                        itmobj.add_reac(reacobj, relref=same_refs_names[0])
                    else:
                        dialog = MultipleChoiceDialog(same_refs_names)
                        dialog.exec()
                        if not dialog.ok:
                            return
                        print("len >1")
                        itmobj.add_reac(reacobj, relref=dialog.items[0])
                else:
                    itmobj.add_reac(reacobj)

    def update_data(self, data=None, filter=None):
        # print("updating sheet", len(self.data.mechs.keys()))
        if data is not None:
            self.data = data
        if filter is None:
            filter = self.parent.combo_filter.currentText()

        self.clear()
        self.setColumnCount(len(self.data.mechs.keys()))
        for col, colname in enumerate(self.data.mechs.keys()):
            # print(col, colname)
            self.setHorizontalHeaderItem(col, QTableWidgetItem(colname))
            header = self.horizontalHeader()
            header.setSectionResizeMode(col, QHeaderView.ResizeToContents)

        # self.setHorizontalHeaderLabels(self.data.mechs.keys())

        maxitm = 0

        # put refs first
        for col, mec in enumerate(self.data.get_mechs()):
            refitems =  mec.get_itms(tp="ref")
            if len(refitems) * 3 > maxitm:
                maxitm = len(refitems) * 3
                self.setRowCount(maxitm)
            row = 0
            for itm in refitems:
                self.setVerticalHeaderItem(row + 0, QTableWidgetItem("Name"))
                self.setVerticalHeaderItem(row + 1, QTableWidgetItem("Energy"))
                self.setVerticalHeaderItem(row + 2, QTableWidgetItem(""))

                font = QFont()
                widget_item = QTableWidgetItem(itm.name)
                font.fromString(itm.name_font)
                widget_item.setFont(font)
                widget_item.setForeground(QColor(120, 120, 250))
                self.setItem(row + 0, col, widget_item)

                font = QFont()
                widget_item = QTableWidgetItem("{:.4f}".format(itm.energy))
                font.fromString(itm.energy_font)
                widget_item.setFont(font)
                self.setItem(row + 1, col, widget_item)

                self.setItem(row + 2, col, QTableWidgetItem(""))
                row += 3
        
        # put the rest
        startrow = maxitm; print(maxitm); 
        
        for col, mec in enumerate(self.data.get_mechs()):
            #mec.get_itms(tp="min") + mec.get_itms(tp="ts") + mec.get_itms(tp=None)
            allitems = mec.get_itms(exclude="ref") 
            #print("all items ", allitems, maxitm)
            if len(allitems) * 3 + startrow > maxitm:
                maxitm = len(allitems) * 3 + startrow
                self.setRowCount(maxitm)
            row = startrow
                # if col == 0:
                # tento que arreglarlo para poner "Name" "Energy" solo cuando haga falta
            for itm in allitems:
                if filter=="None":
                    if itm.tp is not None:
                        continue
                elif filter == "min":
                    if itm.tp != "min":
                        continue
                elif filter == "ts":
                    if itm.tp != "ts":
                        continue
                elif filter == "ref":
                    if itm.tp != "ref":
                        continue
                elif filter == "merged":
                    if not itm.merged:
                        continue
                elif filter == "freqs":
                    if not itm.freqs.vibs:
                        continue
                elif filter == "no freqs":
                    if itm.freqs.vibs:
                        continue
                elif filter == "all":
                    pass
                elif filter != "":
                    if filter not in itm.name:
                        continue

                self.setVerticalHeaderItem(row + 0, QTableWidgetItem("Name"))
                self.setVerticalHeaderItem(row + 1, QTableWidgetItem("Energy"))
                self.setVerticalHeaderItem(row + 2, QTableWidgetItem(""))

                font = QFont()
                widget_item = QTableWidgetItem(itm.name)
                font.fromString(itm.name_font)
                widget_item.setFont(font)
                if itm.tp == "ts":
                    widget_item.setForeground(QColor(250, 120, 120))
                #widget_item.setBackground(QColor(255, 128, 128))
                self.setItem(row + 0, col, widget_item)

                font = QFont()
                widget_item = QTableWidgetItem("{:.4f}".format(itm.energy))
                font.fromString(itm.energy_font)
                widget_item.setFont(font)
                self.setItem(row + 1, col, widget_item)

                self.setItem(row + 2, col, QTableWidgetItem(""))
                row += 3

    def on_itm_click(self, itm):
        self.editing = False
        self.show_itm_info(itm)

    def show_itm_info(self, cellobj):
        col, row = cellobj.column(), cellobj.row()
        colname = self.horizontalHeaderItem(col).text()
        rowobject = self.verticalHeaderItem(row)
        if rowobject is None:
            return
        rowname = rowobject.text()
        if rowname == "Name":
            itmname = cellobj.text()
        elif rowname == "Energy":
            itmname = self.item(row - 1, col).text()
        else:
            itmname = ""
        if not itmname:
            return
        mec = self.data.get_mech(colname)
        itmobj = mec.get_itm(itmname)
        # set combo box temperatures----------------------------
        comboT = self.parent.comboT.currentText()
        self.parent.comboT.clear()
        new_T_list = [str(T) for T in itmobj.temps]
        self.parent.comboT.addItems(new_T_list)
        if comboT in new_T_list:
            self.parent.comboT.setCurrentText(comboT)
        # set combo box temperatures----------------------------
        self.parent.ztable.update_data(itmobj)
        self.parent.reltable.update_data(itmobj)

        li = ""
        li += "UNIT:" + mec.unit + ",  "
        li += "TYPE: " + str(itmobj.tp) + ",  "
        li += "FREQS: " + str(len(itmobj.freqs.vibs)) + " modes, "
        li += "ATOMS: " + ", ".join([str(v) + " "+ str(k) for k, v in itmobj.struct.atoms.items()]) + "\n"
        li += "PG: "+ str(itmobj.pg) + "\n"
        li += "SPIN: "+ str(itmobj.spin) + "\n"
        if itmobj.tp == "mecp":
            li += "MECP PROPS: RedMass: " + str(itmobj.redmass) + ", SOC: " + str(itmobj.SOC)
            li += ", Grad: " + str(itmobj.grad) + ", DiffGrad: " + str(itmobj.diffgrad) + "\n"
        li += "COMMENT: " + itmobj.cm + "\n"
        self.parent.info_label.setText(li)

        if self.parent.openglview:
            self.parent.openGLWidget.update_data(itmobj.struct)
        else:
            # update photo in qlabel self.image
            if itmobj.photos:
                pic = QPixmap(os.environ["HOME"] + "/.datamanager/photos/" + itmobj.photos[0])
                pic = pic.scaled(300, 300)
                self.parent.image.setPixmap(pic)
            else:
                self.parent.image.clear()
        # print(itmobj.thermo)

    def on_itm_dclick(self, itm):
        self.editing = True

class RefSheet(QTableWidget):
    def __init__(self, parent, cols=10, tp="ref"):
        QTableWidget.__init__(self)
        self.parent = parent
        self.tp = tp
        if tp == "reac":
            self.vlabels = [" \u0394E", " \u0394ZPE", " \u0394H", " \u0394G", " \u0394S", " k", " A"]
        else:
            self.vlabels = [" E", " ZPE",  " \u0394H", " G", " \u0394S"]

        self.itmobj = None

        self.setRowCount(len(self.vlabels))
        for row, vl in enumerate(self.vlabels):
            self.setVerticalHeaderItem(row, QTableWidgetItem(vl))
        self.setColumnCount(cols)

        self.customize_corner_button()

        # Bindings -------------------------------------------
        headers = self.horizontalHeader()
        headers.setSectionsClickable(True)
        self.itemClicked.connect(self.on_itm_clicked)
        headers.setContextMenuPolicy(Qt.CustomContextMenu)
        headers.customContextMenuRequested.connect(self.header_popup)
        # Bindings -------------------------------------------

    def customize_corner_button(self):
        if self.tp  == "ref":
          btn = QPushButton("Relative", self)
        else:
          btn = QPushButton("Reaction", self)
        btn.move(0, 0)
        btn.resize(100, self.horizontalHeader().height())
        self.verticalHeader().setFixedWidth(100)
        btn.raise_()

    def header_popup(self, event):
        if self.tp == "ref":
            return
        # self.horizontalHeader().sectionClicked.emit(1)
        menu = QMenu(self)
        gotoitm = menu.addAction("Go to item")
        action = menu.exec_(self.mapToGlobal(event))
        if action == gotoitm:
            self.go_to_itm()

    def go_to_itm(self):
        itmname = self.horizontalHeaderItem(self.currentColumn()).text()

        colname = self.itmobj.mech.name
        mecob = self.itmobj.mech

        row = mecob.get_itm_index(itmname)*3+1
        col = self.parent.data.get_mech_index(colname)
        self.parent.table.clearSelection()
        self.parent.table.setFocus()
        itm = self.parent.table.item(row, col)
        itm.setSelected(True)
        self.parent.table.scrollToItem(itm, QAbstractItemView.PositionAtCenter)

    def on_itm_clicked(self, cellobj):
        col, row = cellobj.column(), cellobj.row()
        colname = self.horizontalHeaderItem(col).text()
        if self.tp == "ref":
            ref = self.itmobj.get_ref(colname)
            if ref is None:
                return
            text = self.parent.info_label.toPlainText()
            text += "\nCoeficients:   " + " ".join([f"{c:.2f}" + " " + str(ref.name) + ",   "
                                           for c, ref in zip(ref.coefs, ref.refs)])
            self.parent.info_label.setText(text)

    def keyPressEvent(self, event):
        key = event.key()
        if key == Qt.Key_Delete:
            self.del_itm()

    def del_itm(self):
        itms = self.selectedIndexes()
        for itm in itms:
            colobj = self.horizontalHeaderItem(itm.column())
            if colobj is None:
                return
            colname = colobj.text()
            if colname == "":
                continue

            if self.tp == "ref":
                self.itmobj.del_ref(colname)
            else:
                self.itmobj.del_reac(colname)

        self.update_data(self.itmobj)

    def contextMenuEvent(self, event):
        menu = QMenu(self)
        tempdep = menu.addAction("Temperature dependency")
        action = menu.exec_(self.mapToGlobal(event.pos()))
        if action == tempdep:
            self.on_temp_dependency()

    def on_temp_dependency(self):
        if self.tp == "ref":
            return
        if self.itmobj.tp != "ts":
            return
        trange, ok = QInputDialog.getText(self, "Temperature range", "Temperature range")
        if not ok:
            return
        if trange:
            if "-" in trange:
                temps = trange.split("-")
            elif ":" in trange:
                temps = trange.split(":")
            else:
                temps = trange.split()
            temps = [float(t.rstrip().lstrip()) for t in temps]
        else:
            temps = []
        reacname = self.horizontalHeaderItem(self.currentColumn()).text()
        reacobj = self.itmobj.get_reac(reacname)
        reacobj.get_kinetic_temp_dependency(temps)

    def update_data(self, itmobj=None):
        if itmobj is not None:
            self.itmobj = itmobj
        if self.itmobj is None:
            return

        try:
            T = float(self.parent.comboT.currentText())
        except ValueError:
            T = 0.0

        if self.tp == "ref":
            zobjlist = self.itmobj.get_refs()
        else:
            zobjlist = self.itmobj.get_reacs()

        self.clearContents()
        self.setColumnCount(6)
        self.setHorizontalHeaderLabels(["", "", "", "", "", ""])
        if not zobjlist:
            for col in range(6):
                header = self.horizontalHeader()
                header.setSectionResizeMode(col, QHeaderView.ResizeToContents)

                for row in range(self.rowCount()):
                    self.setItem(row, col, QTableWidgetItem(""))
            return

        if self.tp == "ref":
            for col, zobj in enumerate(zobjlist):
                if col + 1 > self.columnCount():
                    self.setColumnCount(col)
                    header = self.horizontalHeader()
                    header.setSectionResizeMode(col, QHeaderView.ResizeToContents)

                self.setHorizontalHeaderItem(col, QTableWidgetItem(zobj.name))
                self.setItem(0, col, QTableWidgetItem("{:.4f}".format(zobj.energy)))
                self.setItem(1, col, QTableWidgetItem("{:.4f}".format(zobj.zpe)))
                if T in zobj.thermo["g"]:
                    self.setItem(3, col, QTableWidgetItem("{:.4f}".format(zobj.g(T))))
                    self.setItem(2, col, QTableWidgetItem("{:.4f}".format(zobj.h(T))))
                    self.setItem(4, col, QTableWidgetItem("{:.4f}".format(zobj.stot(T))))

                else:
                    self.setItem(2, col, QTableWidgetItem(""))
                    self.setItem(3, col, QTableWidgetItem(""))
                    self.setItem(4, col, QTableWidgetItem(""))
                    self.setItem(5, col, QTableWidgetItem(""))
                    self.setItem(6, col, QTableWidgetItem(""))
                    self.setItem(7, col, QTableWidgetItem(""))
                    self.setItem(8, col, QTableWidgetItem(""))
        else:
            for col, zobj in enumerate(zobjlist):
                if col + 1 > self.columnCount():
                    self.setColumnCount(col)
                    header = self.horizontalHeader()
                    header.setSectionResizeMode(col, QHeaderView.ResizeToContents)

                self.setHorizontalHeaderItem(col, QTableWidgetItem(zobj.name))
                self.setItem(0, col, QTableWidgetItem("{:.4f}".format(zobj.energy)))
                self.setItem(1, col, QTableWidgetItem("{:.4f}".format(zobj.zpe)))
                if T in zobj.thermo["g"]:
                    self.setItem(2, col, QTableWidgetItem("{:.4f}".format(zobj.h(T))))
                    self.setItem(3, col, QTableWidgetItem("{:.4f}".format(zobj.g(T))))
                    self.setItem(4, col, QTableWidgetItem("{:.4f}".format(zobj.stot(T))))
                    self.setItem(5, col, QTableWidgetItem("{:.3e}".format(zobj.k(T))))
                    #self.setItem(6, col, QTableWidgetItem("{:.3e}".format(zobj.kq(T))))
                    self.setItem(7, col, QTableWidgetItem("{:.3e}".format(zobj.a(T))))
                    #self.setItem(8, col, QTableWidgetItem("{:.4f}".format(zobj.tkey("ae", T))))
                else:
                    self.setItem(2, col, QTableWidgetItem(""))
                    self.setItem(3, col, QTableWidgetItem(""))
                    self.setItem(4, col, QTableWidgetItem(""))
                    self.setItem(5, col, QTableWidgetItem(""))
                    self.setItem(6, col, QTableWidgetItem(""))
                    self.setItem(7, col, QTableWidgetItem(""))
                    self.setItem(8, col, QTableWidgetItem(""))


class MainWindow(QMainWindow):
    def __init__(self, infile):
        QMainWindow.__init__(self)
        self.current_file = infile
        self.data = Data()
        # self.plot = plot(self.data)
        self.plot = MainCanvas(self, self.data)
        self.last_dir = os.getenv("PWD")
        #self.load_data()

        # Widgets initialization------------------------------------------------------
        self.table = MainSheet(self, self.data)

        self.ztable = RefSheet(self)
        self.reltable = RefSheet(self, tp="reac")

        self.info_label = InfoLabel(self)

        self.labelT = QLabel("T (K):")

        # Widgets initialization------------------------------------------------------

        self.title = 'MechaData'
        self.left = 500
        self.top = 500
        self.width = 1000
        self.height = 800
        self.openglview = True
        self.mainlayout = QHBoxLayout()
        self.splitter1 = QSplitter(Qt.Horizontal)
        # self.splitter1.setStretchFactor(1,10)
        self.splitter2 = QSplitter(Qt.Vertical)
        self.splitter3 = QSplitter(Qt.Horizontal)
        self.ref_table_layout = QHBoxLayout()
        self.splitter_widget = QWidget()

        self.image = QLabel(self)
        self.image.setMaximumSize(300, 300)
        # pic = QPixmap("uno.jpg")
        # self.image.setPixmap(pic)
        self.image.setFrameStyle(QFrame.Box)
        self.image.hide()

        # Opengl widget ------------------
        self.openGLWidget = QGLBegin(self)
        # Opengl widget ------------------
        # color
        # self.palete = self.palette()

        self.toolbar = self.addToolBar("toolbar")
        self.comboT = QComboBox(self.toolbar)
        self.filter_label = QLabel("filter by: ")
        self.combo_filter = QComboBox(self.toolbar)

        self.initUI()
        self.load_data()

    def initUI(self):
        self.setWindowTitle(self.title)
        self.setGeometry(self.left, self.top, self.width, self.height)

        # self.palete.setColor(self.backgroundRole(), Qt.black)
        # self.setPalette(self.palete)
        # self.setStyleSheet("background-color:gray;")

        # Menubar ---------------------------------------------------------
        importMec = QAction("&Import Mechanisms From Folder", self)
        importMec.setShortcut("Ctrl+I")
        importMec.setStatusTip("Import mechanisms from folder")
        importMec.triggered.connect(self.import_mec)

        importFromMecFile = QAction("&Import Mechanisms From File", self)
        importFromMecFile.setShortcut("Ctrl+Shift+I")
        importFromMecFile.setStatusTip("Import mechanisms from file")
        importFromMecFile.triggered.connect(self.import_from_mec_file)

        ExportMecExcel = QAction("&Export To Excel", self)
        ExportMecExcel.setShortcut("Ctrl+E")
        ExportMecExcel.setStatusTip("Export to Excel")
        ExportMecExcel.triggered.connect(self.export_to_excel)

        ExportPictures = QAction("&Export Images", self)
        ExportPictures.triggered.connect(self.export_pictures)

        saveAs = QAction("&Save As", self)
        saveAs.setShortcut("Ctrl+Shift+s")
        saveAs.triggered.connect(self.on_save_as)

        save = QAction("&Save", self)
        save.setShortcut("Ctrl+s")
        save.triggered.connect(self.on_save)

        openfile = QAction("&Open File", self)
        openfile.setShortcut("Ctrl+o")
        openfile.triggered.connect(self.open_file)

        closefile = QAction("&Close file", self)
        closefile.setShortcut("Ctrl+w")
        closefile.triggered.connect(self.close_file)

        close = QAction("&Close", self)
        close.setShortcut("Ctrl+q")
        close.triggered.connect(self.quit)

        newMec = QAction("&New Mechanism", self)
        newMec.setShortcut("Ctrl+m")
        newMec.setStatusTip("Create new mechanism")
        newMec.triggered.connect(self.on_new_mech)

        plot_action = QAction("&Show Plot Window", self)
        plot_action.setShortcut("Ctrl+g")
        plot_action.setStatusTip("Show all energy plots")
        plot_action.triggered.connect(self.on_plot)

        # openglview = QAction("&Open GL view", self)
        # openglview.triggered.connect(self.on_openglview)

        # normalview = QAction("&Normal view", self)
        # normalview.triggered.connect(self.on_normalview)

        self.statusBar()

        mainmenu = self.menuBar()
        filemenu = mainmenu.addMenu("&File")
        filemenu.addAction(importMec)
        filemenu.addAction(importFromMecFile)
        filemenu.addAction(save)
        filemenu.addAction(saveAs)
        filemenu.addAction(ExportMecExcel)
        filemenu.addAction(ExportPictures)
        filemenu.addAction(openfile)
        filemenu.addAction(closefile)
        filemenu.addAction(close)

        editmenu = mainmenu.addMenu("&Edit")
        editmenu.addAction(newMec)

        viewmenu = mainmenu.addMenu("&View")
        viewmenu.addAction(plot_action)

        # Menubar ----------------------------------------------------------------

        # TOOL BAR ----------------------------
        self.comboT.setEditable(True)
        self.combo_filter.addItems([
          "all",  "ts", "min", "ref", "None", "freqs", "no freqs",
            "merged", ""
        ])
        self.combo_filter.setEditable(True)
        self.combo_filter.completer().setCaseSensitivity(Qt.CaseSensitive)

        self.toolbar.addWidget(self.labelT)
        self.toolbar.addWidget(self.comboT)
        self.toolbar.addWidget(self.filter_label)
        self.toolbar.addWidget(self.combo_filter)

        # self.toolbar.show()

        # Add box layout, add table to box layout and add box layout to widget
        # self.mainlayout.addWidget(self.splitter1)
        self.splitter1.addWidget(self.table)
        self.splitter1.addWidget(self.splitter2)
        self.splitter2.addWidget(self.info_label)

        label1 = QLabel("Reference table")
        label2 = QLabel("Reference table")
        #panel1 = 

        self.splitter2.addWidget(self.splitter3)
        self.splitter3.addWidget(self.ztable)
        self.splitter3.addWidget(self.reltable)

        self.splitter2.addWidget(self.image)
        self.splitter2.addWidget(self.openGLWidget)
        self.splitter1.setSizes([300, 100])
        self.setCentralWidget(self.splitter1)

        self.bindings()
        self.table.update_data()

        # Show widget
        self.show()

    def bindings(self):
        self.combo_filter.currentTextChanged.connect(self.on_combo_filter)
        self.comboT.currentTextChanged.connect(self.on_combo_T)
        # self.combo_filter.lineEdit().editingFinished.connect(self.on_combo_filter_enter)
        self.info_label.textchanged.connect(self.on_comment_changed)

    # def keyPressEvent(self, event):
    #     # if event.key == Qt.Key_Q and event.modifiers() == Qt.ControlModifier:
    #     if event.key == Qt.Key_Escape:
    #         self.close()

    def quit(self, e):
        d = QMessageBox()
        d.setStandardButtons(QMessageBox.Ok | QMessageBox.Cancel)
        d.setText("Do you really want to close?")
        if d.exec()==1024:
            self.close()

    def on_comment_changed(self, text):
        text = text[text.find("COMMENT:")+8:]
        self.table.edit_comment(text)

    def on_combo_T(self):
        self.ztable.update_data()
        self.reltable.update_data()

    # def on_combo_filter_enter(self):
    #     text = self.combo_filter.currentText()
    #     print(text)
    #     self.combo_filter.addItem(text)
    #     self.combo_filter.setCurrentText(text)
    #     self.table.update_data(filter=text)

    def on_combo_filter(self):
        text = self.combo_filter.currentText()
        self.combo_filter.setItemText(8, text)
        self.combo_filter.setCurrentText(text)
        self.table.update_data(filter=text)

    def load_data(self):
        if not self.current_file:
            self.data = Data()
            self.last_dir = os.getenv("PWD")
            return

        if self.current_file.endswith(".yaml"):
            self.mec_from_yaml_instructions(self.current_file)
            self.current_file = ""
            return

        with open(self.current_file) as f:
            l = f.read()
        data = json.loads(l)
        self.data = Data()
        self.data.dict_to_data(data)
        # self.plot = plot(self.data)
        self.plot = MainCanvas(self, self.data)
        self.last_dir = os.path.basename(self.current_file)
        self.table.update_data(self.data)

    def on_normalview(self):
        self.openGLWidget.hide()
        self.image.show()
        self.openglview = False

    def on_openglview(self):
        self.openGLWidget.show()
        self.openglview = True
        self.image.hide()

    def open_file(self):
        filename = QFileDialog.getOpenFileName()
        if not filename[0]:
            return
        self.current_file = filename[0]
        self.load_data()
        self.table.update_data(self.data)

    def on_save_as(self):
        filename = QFileDialog.getSaveFileName()[0]
        if not filename:
            return
        if not filename.endswith(".json"):
            filename += ".json"
        try:
            with open(filename, "w") as f:
                json.dump(self.data.to_dict(), f, indent=4)
        except:
            msg = QMessageBox()
            msg.setText("File: " + str(self.current_file) + " not saved!")
            msg.exec()
        else:
            self.current_file = filename
            print("File saved correctly")

    def on_save(self):
        if self.current_file is None or self.current_file == "":
            self.on_save_as()
            return
        if os.path.isfile(self.current_file):
            shutil.copy(self.current_file, "./.temp_datamanger")
        try:
            dicdata = self.data.to_dict()
        except Exception as e:
            msg = QMessageBox()
            msg.setText("File: " + str(self.current_file) + " not saved!\nError converting data to dictionary" + str(e))
            msg.exec()
            return
        try:
            dicdata = json.dumps(dicdata, indent=4)
        except Exception as e:
            msg = QMessageBox()
            msg.setText("File: " + str(self.current_file) + " not saved!\nError converting data to json string" + str(e))
            msg.exec()
            return
        try:
            with open(self.current_file, "w") as f:
                f.write(dicdata)
        except Exception as e:
            msg = QMessageBox()
            msg.setText("File: " + str(self.current_file) + " not saved correctly!\n"+str(e))
            msg.exec()
            try:
                shutil.copy("./.temp_datamanager", self.current_file)
                os.rmdir("./.temp_datamanager")
            except Exception as e:
                print("Could not recover temporal data", str(e))
            # raise
        else:
            print("File saved correctly")

    def close_file(self):
        self.current_file = ""
        self.data = Data()
        self.table.update_data(self.data)

    def on_new_mech(self, e):
        dialog = QInputDialog(self, )
        if dialog.exec() == 0:
            return
        text = dialog.textValue()
        self.data.add_m(text)
        self.table.update_data(self.data)

    def import_mec(self, e):
        filedialog = QFileDialog().getExistingDirectory(self,
                                                        "Choose folder", self.last_dir)
        if not filedialog:
            return

        if self.data.mec_from_folder(filedialog) is None:
            return
        self.last_dir = filedialog
        self.table.update_data()

    def import_from_mec_file(self):
        fd = QFileDialog().getOpenFileName(self, "Choose mechanism file", self.last_dir)
        fname = fd[0]
        if not fname:
            return
        if fname.endswith(".yaml"):
            self.mec_from_yaml_instructions(fname)
            return
        
        with open(fname) as f:
            lines = f.read()
        datadict = json.loads(lines)
        data = Data()
        data.dict_to_data(datadict)
        self.data.add_from_data(data)
        self.table.update_data(self.data)

    def mec_from_yaml_instructions(self, fname):
        import yaml
        with open(fname) as f:
            datadict = yaml.safe_load(f)

        for mecname, mecinfo in datadict.items():
            allitems = mecinfo.get("ts_items", []) + mecinfo.get("min_items", []) + mecinfo.get("ref_items", [])
            mecobj = self.data.mec_from_folders(allitems, mecname)
            if mecobj is None:
                continue
            # update types
            for itmname in mecinfo.get("ts_items", []):
                itmobj = mecobj.get_itm(itmname)
                if itmobj is None:
                    continue
                itmobj.tp = "ts"

            for itmname in mecinfo.get("min_items", []):
                itmobj = mecobj.get_itm(itmname)
                if itmobj is None:
                    continue
                itmobj.tp = "min"

            for itmname in mecinfo.get("ref_items", []):
                itmobj = mecobj.get_itm(itmname)
                if itmobj is None:
                    continue
                itmobj.tp = "ref"    

            for itmname, value in mecinfo.get("merged_items", {}).items():
                # values must be a list of str
                mecobj.merge_itms(value, name=itmname)

            # add references
            if not "refs" in mecinfo:
                continue

            for refname, refinfo in mecinfo["refs"].items():
                mecobj.add_ref(refname, refinfo) # refinfo should be a list

            # apply refs
            # print(mecinfo.get("apply_refs_coefs"))
            for refdic in mecinfo.get("apply_refs_coefs", []):
              for refname, refvalue in refdic.items():
                refobjs = mecobj.get_ref(refname)
                if refobjs is None:
                    continue
                for itm_name, coefs in refvalue.items(): # every element in refvalue is a dictionary
                    print("adding ref: ", refname, " to ", itm_name)
                    if len(coefs) != len(refobjs):
                        print("Number of coefficients are not equal to the number of reference items in Reference ", refname)
                        continue
                    itmobj = mecobj.get_itm(itm_name)
                    if itmobj is None:
                        continue
                    if itmobj.tp == "ref": continue
                    try:
                        ref = itmobj.add_ref(refname, refobjs, coefs) 
                        print("added ref: ", ref.name, ref.energy, ref.coefs)
                    except Exception as e:
                        print("Could not add reference ", refname, " to item ", itm_name, str(e), " with coefs ", coefs)

            # apply refs
            #if "apply_refs" in mecinfo:
            for refname, itmlist in mecinfo.get("apply_refs", {}).items():
                refobjs = mecobj.get_ref(refname)
                if refobjs is None:
                    continue
                if isinstance(itmlist, str):
                    tmp_list = []
                    for itm in list(mecobj.itms.keys()):
                        match = re.search(itmlist, itm)
                        if match:
                            tmp_list.append(itm)
                    itmlist = tmp_list
                    
                for itmname in itmlist:
                    itmobj = mecobj.get_itm(itmname)
                    if itmobj is None:
                        continue
                    if itmobj.tp == "ref": continue
                    try:
                        itmobj.add_ref(refname, refobjs) 
                    except Exception as e:
                        print("Could not add reference ", refname, " to item ", itmname, str(e))

            # change units
            unit = mecinfo.get("units") # can be None
            mecobj.chg_unit(unit) # if None chg_unit does nothing


        # make plots
        for mecname, mecinfo in datadict.items():
            mecobj = self.data.get_mech(mecname)
            if mecobj is None:
                continue
            for plotname, pathvalue in mecinfo.get("plots", {}).items():
                # get all items first
                itms = []
                for itmlist in pathvalue.values(): #ignoring the key: pathway1, etc
                    itms += [mecobj.get_itm(itm_name) for itm_name in itmlist if mecobj.get_itm(itm_name) is not None]
       
                # filter out refs and None
                itms = [itm for itm in itms if itm is not None or itm.tp != "ref"]
       
                commonrefs = mecobj.itms_common_refs(itms)
                #print("searching refs for itms:" , itms, len(commonrefs))
                if not commonrefs:
                    msg = QMessageBox()
                    msg.setText("No common reference for selected intermediates to set up a plot")
                    msg.exec()
                    continue
                elif len(commonrefs) == 1:
                    self.data.add_plot(plotname, commonrefs[0], mecname, [], labels=True)
                    #self.plot = MainCanvas(self, self.data)
                    #self.plot.update_data()

                plotobj = self.data.get_plot(plotname)
                # add itms in individual pathways
                for itmlist in pathvalue.values():
                    itms = [mecobj.get_itm(itm_name) for itm_name in itmlist if mecobj.get_itm(itm_name) is not None]
                    plotobj.add_itms(itms)
                    # add reaction reference
                    for n in range(1, len(itms)):
                        itms[n].add_reac(itms[n-1])

        self.table.update_data(self.data)
        self.plot = MainCanvas(self, self.data)
        self.plot.update_data()
        #self.plot.update_data()
        # temporary save 
        with open(".temp_mechasuit.json", "w") as f:
            json.dump(self.data.to_dict(), f, indent=4)

    def export_to_excel(self, e):
        filename, _ = QFileDialog.getSaveFileName(self, "save file", self.last_dir)
        if not filename:
            return
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Energy"
        ws2 = wb.create_sheet("Free energy")

        col = 1
        for mecobj in self.data.get_mechs():
            row = 1
            # aqui deberia merge the cells
            _ = ws.cell(column=col, row=row, value="{0}".format(mecobj.name))
            ws.merge_cells(start_row=row, start_column=col, end_row=row, end_column=col+1)

            _ = ws2.cell(column=col, row=row, value="{0}".format(mecobj.name))
            ws2.merge_cells(start_row=row, start_column=col, end_row=row, end_column=col+1)
            
            row += 1
            row += 1
            # add the references itms first
            for itmobj in mecobj.get_itms(tp="ref"):
                cell = ws.cell(column=col, row=row, value="{0}".format(itmobj.name))
                cell.font = cell.font.copy(bold=True)
                _ = ws.cell(column=col + 1, row=row, value=itmobj.energy)

                cell = ws2.cell(column=col, row=row, value="{0}".format(itmobj.name))
                cell.font = cell.font.copy(bold=True)
                _ = ws2.cell(column=col + 1, row=row, value=itmobj.energy)

                row += 1

            row += 1
            row += 1
            for itmobj in mecobj.get_itms(exclude="ref"):
                cell = ws.cell(column=col, row=row, value="{0}".format(itmobj.name))
                cell.font = cell.font.copy(bold=True)
                _ = ws.cell(column=col+1, row=row, value=itmobj.energy)

                cell = ws2.cell(column=col, row=row, value="{0}".format(itmobj.name))
                cell.font = cell.font.copy(bold=True)
                _ = ws2.cell(column=col+1, row=row, value=itmobj.energy)

                # include the references(relavtive energies)
                row += 1
                for ref in itmobj.get_refs():
                    cell = ws.cell(column=col, row=row, value="{0}".format("Erel "+ref.name))
                    cell.font = cell.font.copy(italic=True, color="ff0000")
                    _ = ws.cell(column=col + 1, row=row, value=ref.energy)

                    cell = ws2.cell(column=col, row=row, value="{0}".format("Grel "+ref.name))
                    cell.font = cell.font.copy(italic=True, color="ff0000")
                    _ = ws2.cell(column=col + 1, row=row, value=ref.g(298.0))

                    row += 1
                # include the reaction energies
                for reac in itmobj.get_reacs():
                    if itmobj.tp == "ts":
                        label = "act"
                    else:
                        label = "reac"
                    cell = ws.cell(column=col, row=row, value="{0}".format("E"+label+"  "+reac.name))
                    cell.font = cell.font.copy(italic=True, color="0066ff")
                    _ = ws.cell(column=col + 1, row=row, value=reac.energy)

                    cell = ws2.cell(column=col, row=row, value="{0}".format("G"+label+"  "+reac.name))
                    cell.font = cell.font.copy(italic=True, color="0066ff")
                    _ = ws2.cell(column=col + 1, row=row, value=reac.g(298.0))

                    row += 1

                row += 1
            col += 2

        wb.save(filename+".xlsx")

    def export_pictures(self, e):
        from pptx import Presentation
        from pptx.util import Inches
        import shutil
        os.mkdir("./.temp")
        os.chdir(".temp")
        print("init png export")
        prs = Presentation()

        for meco in self.data.get_mechs():
            for n, itmo in enumerate(meco.get_itms()):
                gl = QGLBegin(self)
                gl.resize(1000, 1000)
                gl.update_data(itmo.struct)
                # gl.resizeGL(900, 900)
                gl.show()
                gl.hide()

                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                left = top = width = height = Inches(1)

                # incluyo el nombre del mecanimso
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = meco.name
                p.font.bold = True

                # incluyo el nombre del mecanimso
                top = Inches(1.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = itmo.name
                p.font.bold = True

                # img = QImage(gl.width(), gl.height(), QImage.Format_RGB32)
                img = gl.grabFramebuffer()
                img.save(str(n)+"_"+itmo.name+".png")
                top = Inches(3)
                left = Inches(6)
                pic = slide.shapes.add_picture(str(n)+"_"+itmo.name+".png", left, top, height=Inches(3))
        os.chdir("../")
        shutil.rmtree(".temp", ignore_errors=True)
        prs.save("presentation_fotos.pptx")
        print("done exporting png")

    def on_plot(self):
        self.plot.show()
        self.plot.activateWindow()

    def on_del_plot(self):
        plot_names = self.data.get_plots_names()
        print(plot_names)


if __name__ == '__main__':
    app = QApplication(sys.argv)
    # ex = App()
    sys.exit(app.exec_())
